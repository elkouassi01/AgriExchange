"""
VivriMarket_Parcours_Client.pptx
Présentation du parcours client avec images générées par Pillow
"""
import os, math
from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Taille slide 16:9 ──────────────────────────────────────────────────────
SW, SH = 1920, 1080   # pixels pour les images
W  = Inches(13.33)
H  = Inches(7.5)
TMP = "tmp_slides"
os.makedirs(TMP, exist_ok=True)

# ── Palette ────────────────────────────────────────────────────────────────
def rgb(r,g,b): return (r,g,b)
def rgbA(r,g,b,a=255): return (r,g,b,a)

PALETTES = {
    "cover":    [(13,20,14),(16,163,74),(132,204,22)],
    "overview": [(10,15,30),(30,58,138),(56,189,248)],
    "step1":    [(10,30,20),(5,150,105),(52,211,153)],
    "step2":    [(20,10,40),(109,40,217),(167,139,250)],
    "step3":    [(30,20,5),(180,83,9),(251,146,60)],
    "step4":    [(5,20,35),(3,105,161),(56,189,248)],
    "step5":    [(30,5,5),(185,28,28),(248,113,113)],
    "step6":    [(5,25,5),(21,128,61),(74,222,128)],
    "step7":    [(25,15,0),(161,98,7),(251,191,36)],
    "fidelite": [(20,5,25),(124,58,237),(196,181,253)],
    "stats":    [(5,20,30),(14,116,144),(34,211,238)],
    "end":      [(10,15,10),(20,83,45),(74,222,128)],
}


# ── Image generators ───────────────────────────────────────────────────────

def lerp_color(c1, c2, t):
    return tuple(int(c1[i] + (c2[i]-c1[i])*t) for i in range(3))


def make_gradient_bg(key, w=SW, h=SH):
    """Fond dégradé radial + bruit lumineux."""
    c1, c2, c3 = PALETTES[key]
    img = Image.new("RGB", (w, h))
    draw = ImageDraw.Draw(img)
    # vertical gradient base
    for y in range(h):
        t = y / h
        c = lerp_color(c1, c2, t)
        draw.line([(0,y),(w,y)], fill=c)
    # radial glow en haut à droite
    glow = Image.new("RGBA", (w, h), (0,0,0,0))
    gd = ImageDraw.Draw(glow)
    cx, cy = int(w*0.75), int(h*0.22)
    for r in range(600, 0, -8):
        alpha = int(55 * (1 - r/600))
        gd.ellipse([cx-r, cy-r, cx+r, cy+r], fill=(*c3, alpha))
    img = Image.alpha_composite(img.convert("RGBA"), glow).convert("RGB")
    # second glow bas gauche
    glow2 = Image.new("RGBA", (w, h), (0,0,0,0))
    gd2 = ImageDraw.Draw(glow2)
    cx2, cy2 = int(w*0.15), int(h*0.78)
    for r in range(400, 0, -8):
        alpha = int(40 * (1 - r/400))
        gd2.ellipse([cx2-r, cy2-r, cx2+r, cy2+r], fill=(*c2, alpha))
    img = Image.alpha_composite(img.convert("RGBA"), glow2).convert("RGB")
    return img


def add_dots_grid(img, color, spacing=90, dot_r=2, alpha=60):
    overlay = Image.new("RGBA", img.size, (0,0,0,0))
    d = ImageDraw.Draw(overlay)
    for x in range(0, img.width, spacing):
        for y in range(0, img.height, spacing):
            d.ellipse([x-dot_r,y-dot_r,x+dot_r,y+dot_r], fill=(*color, alpha))
    return Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")


def add_diagonal_lines(img, color, gap=60, alpha=25):
    overlay = Image.new("RGBA", img.size, (0,0,0,0))
    d = ImageDraw.Draw(overlay)
    for i in range(-img.height, img.width+img.height, gap):
        d.line([(i,0),(i+img.height, img.height)], fill=(*color, alpha), width=1)
    return Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")


def add_hexagons(img, color, size=80, alpha=30):
    overlay = Image.new("RGBA", img.size, (0,0,0,0))
    d = ImageDraw.Draw(overlay)
    for row in range(0, img.height//size + 2):
        for col in range(0, img.width//size + 2):
            cx = col * size * 1.5
            cy = row * size * 1.73 + (col % 2) * size * 0.87
            pts = []
            for k in range(6):
                angle = math.radians(60*k - 30)
                pts.append((cx + size*0.5*math.cos(angle),
                             cy + size*0.5*math.sin(angle)))
            d.polygon(pts, outline=(*color, alpha), fill=(0,0,0,0))
    return Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")


def add_circles_deco(img, color, positions, alpha=18):
    overlay = Image.new("RGBA", img.size, (0,0,0,0))
    d = ImageDraw.Draw(overlay)
    for cx, cy, r in positions:
        d.ellipse([cx-r,cy-r,cx+r,cy+r], outline=(*color,alpha), width=3)
        d.ellipse([cx-r//2,cy-r//2,cx+r//2,cy+r//2], outline=(*color,alpha//2), width=2)
    return Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")


def make_phone_mockup(img, accent, x, y, w=320, h=560):
    """Dessine un smartphone stylisé sur l'image."""
    d = ImageDraw.Draw(img)
    # Corps
    r = 36
    d.rounded_rectangle([x,y,x+w,y+h], radius=r, fill=(20,20,20), outline=(*accent,200), width=3)
    # Ecran
    d.rounded_rectangle([x+14,y+50,x+w-14,y+h-50], radius=20,
                         fill=lerp_color(accent, (0,0,0), 0.7))
    # Notch
    d.ellipse([x+w//2-28, y+18, x+w//2+28, y+42], fill=(30,30,30))
    # Home button
    d.ellipse([x+w//2-18, y+h-42, x+w//2+18, y+h-14], outline=(*accent,160), width=2)
    # Screen content lines
    lc = lerp_color(accent, (255,255,255), 0.5)
    for i, lw in enumerate([200,150,180,100]):
        ly = y + 90 + i*55
        d.rounded_rectangle([x+28, ly, x+28+lw, ly+18], radius=6, fill=(*lc, 80))
    # App icon grid
    for row in range(3):
        for col in range(3):
            ix = x + 40 + col*80
            iy = y + 330 + row*60
            ic = lerp_color(accent, (50,50,100), 0.3)
            d.rounded_rectangle([ix,iy,ix+50,iy+38], radius=10, fill=ic)


def make_arrow_flow(img, steps, accent, y_center, step_h=100):
    """Dessine un flow d'étapes horizontal."""
    d = ImageDraw.Draw(img)
    total = len(steps)
    bw = (SW - 200) // total
    for i, (label, emoji) in enumerate(steps):
        lx = 100 + i * bw
        cx = lx + bw // 2
        # Background pill
        c_pill = lerp_color(accent, (20,20,30), 0.6)
        d.rounded_rectangle([lx+10, y_center-40, lx+bw-10, y_center+45],
                              radius=20, fill=c_pill, outline=(*accent, 100), width=2)
        # Circle numéro
        d.ellipse([cx-28,y_center-58,cx+28,y_center-2], fill=accent)
        # Arrow
        if i < total-1:
            ax = lx + bw - 12
            d.polygon([(ax,y_center+2),(ax+22,y_center+2),(ax+22-8,y_center-10),(ax+22+8,y_center+2),(ax+22-8,y_center+14)],
                      fill=(*accent, 160))


def save_bg(img, name):
    path = os.path.join(TMP, f"{name}.png")
    img.save(path, "PNG")
    return path


# ── PPTX helpers ──────────────────────────────────────────────────────────

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_full_bg(slide, path):
    pic = slide.shapes.add_picture(path, 0, 0, W, H)
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def R(slide, l,t,w,h, color, alpha_hex="FF"):
    s = slide.shapes.add_shape(1, l,t,w,h)
    s.fill.solid()
    hex_color = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
    s.fill.fore_color.rgb = RGBColor(*color)
    s.line.fill.background()
    return s


def T(slide, text, l,t,w,h, size=20, bold=False, color=(255,255,255),
      align=PP_ALIGN.LEFT, italic=False, spacing=None):
    box = slide.shapes.add_textbox(l,t,w,h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.space_after = Pt(spacing)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor(*color)
    return box


def T2(slide, lines, l,t,w,h, size=16, color=(255,255,255),
       bold_first=False, align=PP_ALIGN.LEFT, line_spacing=4):
    box = slide.shapes.add_textbox(l,t,w,h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = (i == 0 and bold_first)
        r.font.color.rgb = RGBColor(*color)
    return box


def pill(slide, l,t,w,h, bg_color, text, text_color, size=12, bold=True):
    s = slide.shapes.add_shape(1, l,t,w,h)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(*bg_color)
    s.line.fill.background()
    tf = s.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor(*text_color)


def circle_num(slide, num, l,t, size=Inches(0.7), bg=(74,222,128), fg=(0,0,0)):
    s = slide.shapes.add_shape(9, l, t, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(*bg)
    s.line.fill.background()
    tf = s.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(num)
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = RGBColor(*fg)


def divider(slide, accent, y=Inches(1.1), thickness=Inches(0.04)):
    R(slide, 0, y, W, thickness, accent)


def add_image(slide, path, l, t, w, h):
    slide.shapes.add_picture(path, l, t, w, h)


# ── Image builders par slide ───────────────────────────────────────────────

def bg_cover():
    img = make_gradient_bg("cover")
    img = add_hexagons(img, PALETTES["cover"][2], size=120, alpha=22)
    img = add_circles_deco(img, PALETTES["cover"][2],
                           [(SW-280,180,220),(SW-100,500,150),(250,800,180)])
    # Farmer silhouette / field lines
    d = ImageDraw.Draw(img)
    c = PALETTES["cover"][1]
    # Horizon
    d.line([(0, SH-130),(SW, SH-130)], fill=(*c,60), width=2)
    # Crop rows
    for i in range(12):
        x = 80 + i * 140
        for j in range(4):
            yb = SH - 130 + j*30
            d.ellipse([x-15,yb-8,x+15,yb+8], fill=(*PALETTES["cover"][2],40))
    return save_bg(img, "cover")


def bg_overview():
    img = make_gradient_bg("overview")
    img = add_diagonal_lines(img, PALETTES["overview"][2], gap=50, alpha=20)
    img = add_circles_deco(img, PALETTES["overview"][2],
                           [(SW//2,SH//2,600),(SW//2,SH//2,400),(SW//2,SH//2,200)])
    return save_bg(img, "overview")


def bg_step(key, pattern="dots"):
    img = make_gradient_bg(key)
    if pattern == "dots":
        img = add_dots_grid(img, PALETTES[key][2], spacing=70, alpha=50)
    elif pattern == "hex":
        img = add_hexagons(img, PALETTES[key][2], size=90, alpha=25)
    elif pattern == "diag":
        img = add_diagonal_lines(img, PALETTES[key][2], gap=40, alpha=20)
    # Decorative circles
    img = add_circles_deco(img, PALETTES[key][2],
                           [(SW-300, 200, 280),(100, SH-200, 200)])
    return save_bg(img, key)


def make_phone_image(key):
    """Slide image with phone mockup on the right."""
    img = make_gradient_bg(key).convert("RGBA")
    img_base = img.copy()
    d = ImageDraw.Draw(img_base)
    accent = PALETTES[key][2]
    make_phone_mockup(img_base, accent, SW-440, SH//2-280, 320, 560)
    return img_base.convert("RGB")


def bg_with_phone(key):
    img = make_gradient_bg(key)
    img = add_dots_grid(img, PALETTES[key][2], spacing=80, alpha=40)
    d = ImageDraw.Draw(img)
    accent = PALETTES[key][2]
    make_phone_mockup(img, accent, SW-460, SH//2-280, 320, 560)
    return save_bg(img, f"{key}_phone")


def bg_with_cart(key):
    """Slide avec dessin stylisé panier/produit."""
    img = make_gradient_bg(key)
    img = add_hexagons(img, PALETTES[key][2], size=100, alpha=20)
    d = ImageDraw.Draw(img)
    accent = PALETTES[key][2]
    # Cart symbol
    cx, cy = SW-300, SH//2
    # bag shape
    d.rounded_rectangle([cx-120, cy-100, cx+120, cy+140], radius=20,
                         fill=lerp_color(accent,(0,0,0),0.7), outline=(*accent,200), width=4)
    # handle
    d.arc([cx-70, cy-180, cx+70, cy-60], start=200, end=340,
          fill=(*accent,200), width=8)
    # price tag
    d.rounded_rectangle([cx-60,cy-30,cx+60,cy+40], radius=10,
                         fill=lerp_color(accent,(255,255,255),0.1))
    # FCFA text on bag
    for i,lw in enumerate([90,60,75]):
        ly = cy + 50 + i*22
        d.rounded_rectangle([cx-lw//2, ly, cx+lw//2, ly+12], radius=4,
                             fill=(*lerp_color(accent,(255,255,255),0.4), 120))
    return save_bg(img, f"{key}_cart")


def bg_with_payment(key):
    """Slide avec illustration paiement mobile."""
    img = make_gradient_bg(key)
    img = add_diagonal_lines(img, PALETTES[key][2], gap=45, alpha=18)
    d = ImageDraw.Draw(img)
    accent = PALETTES[key][2]
    # Card
    cx, cy = SW-300, SH//2-40
    d.rounded_rectangle([cx-180,cy-100,cx+180,cy+120], radius=24,
                         fill=lerp_color(accent,(0,0,0),0.65), outline=(*accent,180), width=3)
    # Card chip
    d.rounded_rectangle([cx-140,cy-60,cx-60,cy-10], radius=8,
                         fill=lerp_color(accent,(200,200,0),0.5))
    # Lines on card
    for i,lw in enumerate([220,140]):
        d.rounded_rectangle([cx-140,cy+10+i*35,cx-140+lw,cy+25+i*35], radius=4,
                             fill=(*lerp_color(accent,(255,255,255),0.4),150))
    # Wifi-pay icon above
    for r in [80,55,30]:
        d.arc([cx-r,cy-185-r,cx+r,cy-185+r], start=220, end=320,
              fill=(*accent,int(180*(1-r/100))), width=5)
    d.ellipse([cx-8,cy-185-8,cx+8,cy-185+8], fill=accent)
    # Amount badge
    d.rounded_rectangle([cx-80,cy+145,cx+80,cy+195], radius=20,
                         fill=accent)
    return save_bg(img, f"{key}_pay")


def bg_with_handshake(key):
    """Slide avec illustration contact / handshake."""
    img = make_gradient_bg(key)
    img = add_dots_grid(img, PALETTES[key][2], spacing=75, alpha=45)
    d = ImageDraw.Draw(img)
    accent = PALETTES[key][2]
    # Two speech bubbles
    cx, cy = SW-280, SH//2
    # Bubble 1 (left - buyer)
    b1c = lerp_color(accent,(0,0,0),0.5)
    d.rounded_rectangle([cx-230,cy-160,cx-30,cy-60], radius=18,
                         fill=b1c, outline=(*accent,150), width=2)
    d.polygon([(cx-190,cy-60),(cx-170,cy-30),(cx-150,cy-60)], fill=b1c)
    # Bubble 2 (right - seller)
    b2c = lerp_color(PALETTES[key][1],(255,255,255),0.1)
    d.rounded_rectangle([cx-20,cy-60,cx+200,cy+40], radius=18,
                         fill=b2c, outline=(*accent,120), width=2)
    d.polygon([(cx+30,cy+40),(cx+50,cy+70),(cx+70,cy+40)], fill=b2c)
    # Phone icon
    d.ellipse([cx-30,cy+80,cx+30,cy+140], fill=accent)
    d.ellipse([cx-20,cy+90,cx+20,cy+130], fill=lerp_color(accent,(0,0,0),0.6))
    # Signal waves
    for r in [60,45,30]:
        d.arc([cx-r,cy+60-r//2+110-r//2,cx+r,cy+60+r//2+110+r//2],
              start=200,end=340, fill=(*accent,int(160*(1-r/70))), width=4)
    return save_bg(img, f"{key}_chat")


def bg_stars(key):
    """Slide fidélité avec étoiles."""
    img = make_gradient_bg(key)
    img = add_hexagons(img, PALETTES[key][2], size=110, alpha=18)
    d = ImageDraw.Draw(img)
    accent = PALETTES[key][2]

    def star(d, cx, cy, r_out, r_in, color, alpha=200):
        pts = []
        for k in range(10):
            angle = math.radians(k * 36 - 90)
            r = r_out if k % 2 == 0 else r_in
            pts.append((cx + r*math.cos(angle), cy + r*math.sin(angle)))
        d.polygon(pts, fill=(*color, alpha))

    positions = [(SW-320,SH//2-20,90,40),(SW-200,SH//2-120,65,28),
                 (SW-440,SH//2-80,55,24),(SW-250,SH//2+110,50,22),
                 (SW-170,SH//2+60,40,18)]
    for i,(cx,cy,ro,ri) in enumerate(positions):
        alpha_v = 220 - i*25
        star(d, cx, cy, ro, ri, accent, alpha_v)

    return save_bg(img, f"{key}_stars")


def bg_stats_img(key):
    """Slide stats avec graphes simplifiés."""
    img = make_gradient_bg(key)
    img = add_diagonal_lines(img, PALETTES[key][2], gap=55, alpha=15)
    d = ImageDraw.Draw(img)
    accent = PALETTES[key][2]
    c2 = PALETTES[key][1]
    # Bar chart
    bars = [0.4, 0.65, 0.55, 0.85, 0.75, 0.9]
    bw, gap = 55, 20
    bx = SW - 480
    by = SH//2 + 120
    for i, v in enumerate(bars):
        bh = int(v * 280)
        lx = bx + i*(bw+gap)
        bc = lerp_color(accent, c2, i/len(bars))
        d.rounded_rectangle([lx, by-bh, lx+bw, by], radius=8, fill=(*bc,200))
    # Axis
    d.line([(bx-10, by),(bx + len(bars)*(bw+gap), by)], fill=(*accent,120), width=2)
    # Donut circle
    dcx, dcy, dr = SW-340, SH//2-130, 110
    d.ellipse([dcx-dr,dcy-dr,dcx+dr,dcy+dr], fill=lerp_color(accent,(0,0,0),0.8))
    d.ellipse([dcx-70,dcy-70,dcx+70,dcy+70], fill=lerp_color(accent,(0,0,30),0.7))
    d.arc([dcx-dr,dcy-dr,dcx+dr,dcy+dr], start=-90, end=200,
          fill=accent, width=18)
    return save_bg(img, f"{key}_stats")


def bg_end():
    img = make_gradient_bg("end")
    img = add_hexagons(img, PALETTES["end"][2], size=130, alpha=20)
    d = ImageDraw.Draw(img)
    accent = PALETTES["end"][2]
    # Big checkmark
    cx, cy = SW-300, SH//2
    d.ellipse([cx-160,cy-160,cx+160,cy+160], fill=lerp_color(accent,(0,0,0),0.65),
              outline=(*accent,200), width=5)
    d.line([(cx-80,cy),(cx-20,cy+70),(cx+90,cy-80)], fill=(*accent,230), width=14)
    return save_bg(img, "end")


# ═══════════════════════════════════════════════════════════════════════════════
# Génération des images
# ═══════════════════════════════════════════════════════════════════════════════
print("Generating images...")
paths = {
    "cover":    bg_cover(),
    "overview": bg_overview(),
    "step1":    bg_step("step1","hex"),
    "step2":    bg_with_phone("step2"),
    "step3":    bg_with_cart("step3"),
    "step4":    bg_step("step4","diag"),
    "step5":    bg_with_payment("step5"),
    "step6":    bg_with_handshake("step6"),
    "step7":    bg_step("step7","dots"),
    "fidelite": bg_stars("fidelite"),
    "stats":    bg_stats_img("stats"),
    "end":      bg_end(),
}
print("  Images OK")


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX
# ═══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

GREEN   = (74,222,128)
LIME    = (132,204,22)
BLUE    = (56,189,248)
VIOLET  = (167,139,250)
ORANGE  = (251,146,60)
YELLOW  = (251,191,36)
PINK    = (244,114,182)
CYAN    = (34,211,238)
RED     = (248,113,113)
TEAL    = (45,212,191)
WHITE   = (255,255,255)
DIM     = (113,113,122)
BLACK   = (10,10,10)


def overlay_panel(slide, l,t,w,h, color, opacity=180):
    """Semi-transparent panel via solid fill (approximation)."""
    s = slide.shapes.add_shape(1,l,t,w,h)
    s.fill.solid()
    r,g,b = color
    # darken by opacity ratio
    factor = opacity/255
    s.fill.fore_color.rgb = RGBColor(int(r*factor),int(g*factor),int(b*factor))
    s.line.fill.background()
    return s


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["cover"])

# Left dark overlay
overlay_panel(s, 0,0, Inches(7.2), H, (5,12,8), 210)

# Accent bar
R(s, 0,0, Inches(0.1), H, GREEN)

# Logo badge
sp = s.shapes.add_shape(9, Inches(0.5), Inches(0.6), Inches(1.3), Inches(1.3))
sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(*GREEN); sp.line.fill.background()
tf = sp.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
r = tf.paragraphs[0].add_run(); r.text = "VM"
r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = RGBColor(*BLACK)

T(s, "VivriMarket", Inches(2.1), Inches(0.6), Inches(6), Inches(0.85),
  size=14, color=GREEN, bold=True)

T(s, "Parcours Client", Inches(0.5), Inches(1.85), Inches(7), Inches(1.4),
  size=52, bold=True, color=WHITE)

T(s, "De la découverte à l'achat", Inches(0.5), Inches(3.25), Inches(7), Inches(0.7),
  size=22, color=GREEN)

T(s, "Comment un acheteur trouve un agriculteur\net conclut une transaction sur VivriMarket",
  Inches(0.5), Inches(4.1), Inches(6.8), Inches(1.2),
  size=15, color=(180,200,185))

# 3 stat badges
stats_data = [("8", "étapes clés"), ("300", "FCFA contact"), ("∞", "vendeurs")]
for i, (val, lbl) in enumerate(stats_data):
    lx = Inches(0.5) + i * Inches(2.15)
    R(s, lx, Inches(5.7), Inches(1.9), Inches(1.1), (10,40,20))
    R(s, lx, Inches(5.7), Inches(0.07), Inches(1.1), GREEN)
    T(s, val, lx+Inches(0.2), Inches(5.75), Inches(1.6), Inches(0.55),
      size=26, bold=True, color=GREEN)
    T(s, lbl, lx+Inches(0.2), Inches(6.25), Inches(1.6), Inches(0.4),
      size=11, color=(150,200,160))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — VUE D'ENSEMBLE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["overview"])

overlay_panel(s, 0,0, W, Inches(1.1), (5,10,25), 200)

T(s, "Vue d'ensemble — Le parcours en 8 étapes",
  Inches(0.6), Inches(0.2), Inches(12), Inches(0.75),
  size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Timeline horizontal
steps_ov = [
    ("01","Découverte",       BLUE),
    ("02","Inscription",      VIOLET),
    ("03","Explorer",         ORANGE),
    ("04","Choisir",          CYAN),
    ("05","Débloquer",        RED),
    ("06","Contacter",        GREEN),
    ("07","Acheter",          YELLOW),
    ("08","Fidélité",         PINK),
]

box_w = Inches(1.48)
box_h = Inches(4.5)
gap   = Inches(0.08)
start_x = Inches(0.3)

for i, (num, label, color) in enumerate(steps_ov):
    lx = start_x + i * (box_w + gap)
    # Vertical bar
    R(s, lx, Inches(1.2), box_w, box_h, (10,15,35))
    R(s, lx, Inches(1.2), box_w, Inches(0.07), color)

    # Number circle
    cx_off = lx + box_w/2 - Inches(0.38)
    sp2 = s.shapes.add_shape(9, cx_off, Inches(1.45), Inches(0.76), Inches(0.76))
    sp2.fill.solid(); sp2.fill.fore_color.rgb = RGBColor(*color); sp2.line.fill.background()
    tf2 = sp2.text_frame; tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    r2 = tf2.paragraphs[0].add_run(); r2.text = num
    r2.font.size = Pt(15); r2.font.bold = True; r2.font.color.rgb = RGBColor(*BLACK)

    # Label
    T(s, label, lx+Inches(0.08), Inches(2.45), box_w-Inches(0.1), Inches(0.6),
      size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Arrow between boxes
    if i < len(steps_ov)-1:
        ax = lx + box_w + Inches(0.01)
        T(s, "›", ax, Inches(3.0), gap+Inches(0.05), Inches(0.5),
          size=20, bold=True, color=(80,100,150), align=PP_ALIGN.CENTER)

# Bottom description
descs = [
    "Réseau\nsocial\nfarm",
    "WhatsApp\nOTP\nrapide",
    "Catalogue\nproduits\nlocaux",
    "Profil\nvendeur\ndetaillé",
    "Paiement\nmobile\n300 FCFA",
    "Contact\ndirect\nvendeur",
    "Accord\nprix &\nlivraison",
    "Historique\net\nrecommandation",
]
for i, desc in enumerate(descs):
    lx = start_x + i * (box_w + gap)
    T(s, desc, lx+Inches(0.05), Inches(3.2), box_w-Inches(0.05), Inches(2.3),
      size=10, color=(140,160,190), align=PP_ALIGN.CENTER)

T(s, "Chaque étape est pensée pour connecter acheteurs et agriculteurs locaux de façon simple et sécurisée.",
  Inches(0.4), Inches(5.9), Inches(12.5), Inches(0.45),
  size=12, italic=True, color=(80,100,130), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — ÉTAPE 1 : DÉCOUVERTE
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["step1"])

overlay_panel(s, 0,0, Inches(7.5), H, (5,18,10), 200)

R(s, 0,0, Inches(0.12), H, GREEN)

circle_num(s, "01", Inches(0.35), Inches(0.35), Inches(0.8), GREEN, BLACK)
T(s, "Étape 1  —  Découverte", Inches(1.35), Inches(0.38), Inches(8), Inches(0.65),
  size=24, bold=True, color=GREEN)
R(s, Inches(0.35), Inches(1.15), Inches(7.1), Inches(0.04), GREEN)

T(s, "Comment les acheteurs\ntrouvent VivriMarket",
  Inches(0.4), Inches(1.35), Inches(7), Inches(1.3),
  size=30, bold=True, color=WHITE)

T(s, "L'acheteur entend parler de VivriMarket via son réseau, les réseaux sociaux ou un marché local.",
  Inches(0.4), Inches(2.75), Inches(7), Inches(0.75),
  size=15, color=(180,220,190))

channels = [
    ("WhatsApp", "Partage entre amis et famille",     (37,211,102)),
    ("Facebook", "Publications de fermiers locaux",   (66,103,178)),
    ("Bouche à oreille", "Marchés et communautés",    YELLOW),
    ("Google", "Recherche de produits locaux",         RED),
]
for i, (ch, desc, color) in enumerate(channels):
    col, row = i%2, i//2
    lx = Inches(0.4) + col*Inches(3.5)
    ty = Inches(3.75) + row*Inches(1.45)
    R(s, lx, ty, Inches(3.2), Inches(1.2), (8,30,15))
    R(s, lx, ty, Inches(0.12), Inches(1.2), color)
    T(s, ch, lx+Inches(0.25), ty+Inches(0.1), Inches(2.8), Inches(0.5),
      size=15, bold=True, color=WHITE)
    T(s, desc, lx+Inches(0.25), ty+Inches(0.6), Inches(2.8), Inches(0.45),
      size=11, color=(160,200,170))

T(s, '"La premiere plateforme qui\nconnecte vraiment les agriculteurs\nde notre village aux acheteurs."',
  Inches(0.4), Inches(6.6), Inches(7), Inches(0.7),
  size=11, italic=True, color=(100,150,110))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — ÉTAPE 2 : INSCRIPTION
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["step2"])

overlay_panel(s, 0,0, Inches(7.8), H, (15,5,30), 210)

R(s, 0,0, Inches(0.12), H, VIOLET)

circle_num(s, "02", Inches(0.35), Inches(0.35), Inches(0.8), VIOLET, BLACK)
T(s, "Étape 2  —  Inscription", Inches(1.35), Inches(0.38), Inches(8), Inches(0.65),
  size=24, bold=True, color=VIOLET)
R(s, Inches(0.35), Inches(1.15), Inches(7.3), Inches(0.04), VIOLET)

T(s, "Créer un compte\nen quelques secondes",
  Inches(0.4), Inches(1.35), Inches(7.2), Inches(1.3),
  size=30, bold=True, color=WHITE)

# Form mockup
fields = [
    ("Nom complet",    "ex: Kouassi Aya"),
    ("Email",         "ex: aya@email.com"),
    ("Telephone",     "ex: +225 07 XX XX XX"),
    ("Mot de passe",  "••••••••••"),
]
for i, (label, placeholder) in enumerate(fields):
    ty = Inches(2.9) + i * Inches(0.88)
    T(s, label, Inches(0.5), ty, Inches(6.5), Inches(0.32),
      size=10, bold=True, color=VIOLET)
    R(s, Inches(0.5), ty+Inches(0.33), Inches(6.5), Inches(0.48),
      (30,15,55))
    R(s, Inches(0.5), ty+Inches(0.33), Inches(0.08), Inches(0.48), VIOLET)
    T(s, placeholder, Inches(0.72), ty+Inches(0.38), Inches(6.2), Inches(0.38),
      size=12, color=(120,100,160))

# OTP badge
R(s, Inches(0.5), Inches(6.65), Inches(6.5), Inches(0.6), (60,20,100))
T(s, "  Verification OTP par WhatsApp — code valable 10 minutes",
  Inches(0.5), Inches(6.65), Inches(6.5), Inches(0.6),
  size=12, bold=True, color=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — ÉTAPE 3 : EXPLORER LES PRODUITS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["step3"])

overlay_panel(s, 0,0, Inches(7.8), H, (25,12,0), 205)

R(s, 0,0, Inches(0.12), H, ORANGE)

circle_num(s, "03", Inches(0.35), Inches(0.35), Inches(0.8), ORANGE, BLACK)
T(s, "Étape 3  —  Explorer les produits", Inches(1.35), Inches(0.38), Inches(8), Inches(0.65),
  size=24, bold=True, color=ORANGE)
R(s, Inches(0.35), Inches(1.15), Inches(7.3), Inches(0.04), ORANGE)

T(s, "Un marché agricole\ntout dans la poche",
  Inches(0.4), Inches(1.35), Inches(7.2), Inches(1.3),
  size=30, bold=True, color=WHITE)

categories = [
    ("Cereales & Legumes",   "Mais, riz, tomates, ignames..."),
    ("Fruits tropicaux",      "Ananas, mangues, papayes..."),
    ("Elevage & Laitier",    "Poulets, oeufs, lait frais..."),
    ("Cultures de rente",    "Cacao, cafe, anacarde..."),
]
for i, (cat, items) in enumerate(categories):
    col, row = i%2, i//2
    lx = Inches(0.4) + col*Inches(3.6)
    ty = Inches(2.8) + row*Inches(1.5)
    R(s, lx, ty, Inches(3.3), Inches(1.25), (35,18,5))
    R(s, lx, ty, Inches(0.12), Inches(1.25), ORANGE)
    T(s, cat, lx+Inches(0.25), ty+Inches(0.12),
      Inches(2.95), Inches(0.5), size=14, bold=True, color=WHITE)
    T(s, items, lx+Inches(0.25), ty+Inches(0.65),
      Inches(2.95), Inches(0.45), size=11, color=(200,160,100))

T(s, "Filtres disponibles : localisation, prix, disponibilite, categorie",
  Inches(0.4), Inches(6.2), Inches(7.2), Inches(0.45),
  size=12, color=(200,160,80))
pill(s, Inches(0.4), Inches(6.75), Inches(7.2), Inches(0.5),
    ORANGE, "  Acces 100% GRATUIT — Parcourir le catalogue ne coute rien  ",
    BLACK, size=12)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — ÉTAPE 4 : CHOISIR UN VENDEUR
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["step4"])

overlay_panel(s, 0,0, Inches(7.8), H, (2,18,30), 205)

R(s, 0,0, Inches(0.12), H, CYAN)

circle_num(s, "04", Inches(0.35), Inches(0.35), Inches(0.8), CYAN, BLACK)
T(s, "Étape 4  —  Choisir un vendeur", Inches(1.35), Inches(0.38), Inches(8), Inches(0.65),
  size=24, bold=True, color=CYAN)
R(s, Inches(0.35), Inches(1.15), Inches(7.3), Inches(0.04), CYAN)

T(s, "Consulter la fiche\ndu produit & du vendeur",
  Inches(0.4), Inches(1.35), Inches(7.2), Inches(1.3),
  size=30, bold=True, color=WHITE)

# Product card mockup
R(s, Inches(0.4), Inches(2.85), Inches(7.1), Inches(4.0), (5,20,35))
R(s, Inches(0.4), Inches(2.85), W*0+Inches(7.1), Inches(0.07), CYAN)

card_items = [
    ("Produit",   "Tomates fraîches — 10 kg"),
    ("Vendeur",   "Koffi Jean-Baptiste"),
    ("Localisation","Yamoussoukro, Côte d'Ivoire"),
    ("Prix",      "5 000 FCFA / sac"),
    ("Note",      "4.8 / 5  (32 avis)"),
    ("Statut",    "Disponible immediatement"),
]
for i, (lbl, val) in enumerate(card_items):
    ty = Inches(3.05) + i*Inches(0.58)
    T(s, lbl, Inches(0.65), ty, Inches(2.2), Inches(0.48),
      size=11, bold=True, color=(100,170,200))
    T(s, val, Inches(2.95), ty, Inches(4.4), Inches(0.48),
      size=12, color=WHITE)
    if i < len(card_items)-1:
        R(s, Inches(0.55), ty+Inches(0.52), Inches(6.8), Inches(0.02), (20,40,60))

pill(s, Inches(0.65), Inches(6.55), Inches(3.1), Inches(0.55),
    CYAN, "Voir le contact", BLACK, size=13)
pill(s, Inches(3.9), Inches(6.55), Inches(3.4), Inches(0.55),
    (20,50,80), "Contacter sur WhatsApp", CYAN, size=13)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — ÉTAPE 5 : DÉBLOQUER LE CONTACT
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["step5"])

overlay_panel(s, 0,0, Inches(7.8), H, (28,5,5), 205)

R(s, 0,0, Inches(0.12), H, RED)

circle_num(s, "05", Inches(0.35), Inches(0.35), Inches(0.8), RED, BLACK)
T(s, "Étape 5  —  Débloquer le contact", Inches(1.35), Inches(0.38), Inches(8), Inches(0.65),
  size=24, bold=True, color=RED)
R(s, Inches(0.35), Inches(1.15), Inches(7.3), Inches(0.04), RED)

T(s, "Un micro-paiement de\n300 FCFA pour accéder\nau vendeur",
  Inches(0.4), Inches(1.35), Inches(7.2), Inches(1.7),
  size=28, bold=True, color=WHITE)

# Payment steps
pay_steps = [
    ("1", "Cliquer «Débloquer le contact»",                     RED),
    ("2", "Choisir son moyen de paiement\n(Orange Money, MTN, Wave...)", ORANGE),
    ("3", "Valider la transaction de 300 FCFA",                 YELLOW),
    ("4", "Le numéro du vendeur apparaît instantanément",       GREEN),
]
for i, (num, desc, color) in enumerate(pay_steps):
    ty = Inches(3.25) + i * Inches(0.95)
    R(s, Inches(0.4), ty, Inches(7.1), Inches(0.83), (35,8,8))
    R(s, Inches(0.4), ty, Inches(0.12), Inches(0.83), color)
    sp = s.shapes.add_shape(9, Inches(0.65), ty+Inches(0.15), Inches(0.55), Inches(0.55))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(*color); sp.line.fill.background()
    tf = sp.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = num
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = RGBColor(*BLACK)
    T(s, desc, Inches(1.35), ty+Inches(0.15), Inches(5.95), Inches(0.65),
      size=13, color=WHITE)

R(s, Inches(0.4), Inches(7.05), Inches(7.1), Inches(0.25), (50,10,10))
T(s, "  Contact valable 30 jours — remboursé si le vendeur ne repond pas",
  Inches(0.4), Inches(7.05), Inches(7.1), Inches(0.35),
  size=11, color=(255,150,150))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — ÉTAPE 6 : CONTACTER LE VENDEUR
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["step6"])

overlay_panel(s, 0,0, Inches(7.8), H, (2,20,5), 200)

R(s, 0,0, Inches(0.12), H, GREEN)

circle_num(s, "06", Inches(0.35), Inches(0.35), Inches(0.8), GREEN, BLACK)
T(s, "Étape 6  —  Contacter le vendeur", Inches(1.35), Inches(0.38), Inches(8), Inches(0.65),
  size=24, bold=True, color=GREEN)
R(s, Inches(0.35), Inches(1.15), Inches(7.3), Inches(0.04), GREEN)

T(s, "Le contact direct\nvendeur–acheteur",
  Inches(0.4), Inches(1.35), Inches(7.2), Inches(1.3),
  size=32, bold=True, color=WHITE)

T(s, "Une fois le contact débloqué, l'acheteur dispose\ndu numéro WhatsApp / téléphone du vendeur.",
  Inches(0.4), Inches(2.75), Inches(7.2), Inches(0.9),
  size=14, color=(160,220,170))

channels2 = [
    ("Appel direct",    "Negocier le prix et la quantite",      GREEN),
    ("WhatsApp",        "Envoyer des photos du besoin",          (37,211,102)),
    ("SMS",             "Confirmer la commande par ecrit",       CYAN),
    ("Rencontre",       "Visiter l'exploitation si besoin",      YELLOW),
]
for i, (mode, desc, color) in enumerate(channels2):
    col, row = i%2, i//2
    lx = Inches(0.4) + col*Inches(3.6)
    ty = Inches(3.85) + row*Inches(1.5)
    R(s, lx, ty, Inches(3.3), Inches(1.25), (5,25,10))
    R(s, lx, ty, Inches(0.12), Inches(1.25), color)
    T(s, mode, lx+Inches(0.25), ty+Inches(0.12), Inches(2.95), Inches(0.5),
      size=15, bold=True, color=WHITE)
    T(s, desc, lx+Inches(0.25), ty+Inches(0.65), Inches(2.95), Inches(0.45),
      size=11, color=(140,200,150))

T(s, "La plateforme garantit la visibilite du contact pendant 30 jours.",
  Inches(0.4), Inches(7.05), Inches(7.2), Inches(0.35),
  size=11, italic=True, color=(80,150,100))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — ÉTAPE 7 : CONCLURE L'ACHAT
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["step7"])

overlay_panel(s, 0,0, Inches(7.8), H, (22,14,0), 205)

R(s, 0,0, Inches(0.12), H, YELLOW)

circle_num(s, "07", Inches(0.35), Inches(0.35), Inches(0.8), YELLOW, BLACK)
T(s, "Étape 7  —  Conclure l'achat", Inches(1.35), Inches(0.38), Inches(8), Inches(0.65),
  size=24, bold=True, color=YELLOW)
R(s, Inches(0.35), Inches(1.15), Inches(7.3), Inches(0.04), YELLOW)

T(s, "De l'accord à\nla livraison",
  Inches(0.4), Inches(1.35), Inches(7.2), Inches(1.2),
  size=32, bold=True, color=WHITE)

T(s, "L'acheteur et le vendeur se mettent d'accord sur le prix, la quantité et les modalités de livraison.",
  Inches(0.4), Inches(2.65), Inches(7.2), Inches(0.75),
  size=14, color=(220,190,100))

deal_steps = [
    ("Negociation",  "Prix et quantite discutes directement",  YELLOW),
    ("Commande",     "Accord confirme par WhatsApp ou SMS",    ORANGE),
    ("Livraison",    "Transport arrange par les deux parties", RED),
    ("Reception",    "Acheteur confirme la reception",        GREEN),
]
for i, (title, desc, color) in enumerate(deal_steps):
    lx = Inches(0.4) + (i % 2) * Inches(3.6)
    ty = Inches(3.6) + (i // 2) * Inches(1.5)
    R(s, lx, ty, Inches(3.3), Inches(1.25), (30,20,5))
    R(s, lx, ty, Inches(0.12), Inches(1.25), color)
    T(s, title, lx+Inches(0.25), ty+Inches(0.12), Inches(2.95), Inches(0.5),
      size=15, bold=True, color=WHITE)
    T(s, desc, lx+Inches(0.25), ty+Inches(0.65), Inches(2.95), Inches(0.45),
      size=11, color=(200,170,100))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — ÉTAPE 8 : FIDÉLISATION
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["fidelite"])

overlay_panel(s, 0,0, Inches(7.8), H, (18,5,22), 205)

R(s, 0,0, Inches(0.12), H, VIOLET)

circle_num(s, "08", Inches(0.35), Inches(0.35), Inches(0.8), VIOLET, BLACK)
T(s, "Étape 8  —  Fidélisation", Inches(1.35), Inches(0.38), Inches(8), Inches(0.65),
  size=24, bold=True, color=VIOLET)
R(s, Inches(0.35), Inches(1.15), Inches(7.3), Inches(0.04), VIOLET)

T(s, "Revenir, noter,\nrecommander",
  Inches(0.4), Inches(1.35), Inches(7.2), Inches(1.25),
  size=32, bold=True, color=WHITE)

T(s, "Après l'achat, l'acheteur devient un ambassadeur de la plateforme.",
  Inches(0.4), Inches(2.7), Inches(7.2), Inches(0.55),
  size=14, color=(180,150,230))

loyalty = [
    ("Laisser un avis",       "Note et commentaire visible sur le profil vendeur",   VIOLET),
    ("Historique d'achats",   "Retrouver ses vendeurs habituels facilement",          PINK),
    ("Partager",              "Inviter des amis agriculteurs ou acheteurs",           CYAN),
    ("Programme fidélité",    "Réductions sur les prochains déblocages",              YELLOW),
]
for i, (title, desc, color) in enumerate(loyalty):
    ty = Inches(3.5) + i * Inches(0.85)
    R(s, Inches(0.4), ty, Inches(7.1), Inches(0.73), (22,8,28))
    R(s, Inches(0.4), ty, Inches(0.12), Inches(0.73), color)
    T(s, title, Inches(0.65), ty+Inches(0.08), Inches(2.3), Inches(0.55),
      size=13, bold=True, color=WHITE)
    T(s, desc, Inches(3.0), ty+Inches(0.12), Inches(4.35), Inches(0.5),
      size=12, color=(170,140,210))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — CHIFFRES CLÉS
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["stats"])

overlay_panel(s, 0,0, W, Inches(1.2), (2,12,22), 220)

T(s, "Chiffres clés du parcours client",
  Inches(0.6), Inches(0.2), Inches(12.1), Inches(0.8),
  size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
R(s, Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.05), CYAN)

big_stats = [
    ("300 FCFA",     "Coût pour débloquer\nun contact vendeur",    CYAN),
    ("30 jours",     "Durée de validité\nd'un contact débloqué",   BLUE),
    ("< 1 min",      "Temps moyen pour\ntrouver un vendeur",       GREEN),
    ("100%",         "Catalogue produits\naccès gratuit",          LIME),
    ("4.8/5",        "Note moyenne\ndes vendeurs actifs",          YELLOW),
    ("Remboursé",    "Si le vendeur\nne répond pas",               RED),
]

for i, (val, desc, color) in enumerate(big_stats):
    col = i % 3
    row = i // 3
    lx = Inches(0.6) + col * Inches(4.1)
    ty = Inches(1.35) + row * Inches(2.7)
    R(s, lx, ty, Inches(3.85), Inches(2.45), (5,18,30))
    R(s, lx, ty, Inches(3.85), Inches(0.08), color)
    T(s, val, lx, ty+Inches(0.25), Inches(3.85), Inches(1.0),
      size=38, bold=True, color=color, align=PP_ALIGN.CENTER)
    T(s, desc, lx, ty+Inches(1.3), Inches(3.85), Inches(0.9),
      size=13, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — FIN
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
set_full_bg(s, paths["end"])

overlay_panel(s, 0,0, Inches(8.0), H, (5,15,8), 200)

R(s, 0,0, Inches(0.12), H, GREEN)
R(s, 0, H-Inches(0.12), W, Inches(0.12), GREEN)

T(s, "Un parcours simple,\ntransparent et local.",
  Inches(0.5), Inches(1.2), Inches(7.5), Inches(1.9),
  size=38, bold=True, color=WHITE)

T(s, "VivriMarket connecte directement acheteurs\net agriculteurs locaux, sans intermediaire,\npour des achats alimentaires plus frais,\nplus rapides et plus equitables.",
  Inches(0.5), Inches(3.2), Inches(7.3), Inches(2.0),
  size=16, color=(160,220,175))

# Key points
for i, (pt, color) in enumerate([
    ("Inscription gratuite",              GREEN),
    ("Catalogue 100% gratuit",            LIME),
    ("300 FCFA seulement pour le contact", CYAN),
    ("Contact garanti ou rembourse",       YELLOW),
]):
    ty = Inches(5.5) + i * Inches(0.42)
    sp = s.shapes.add_shape(9, Inches(0.5), ty+Inches(0.06), Inches(0.28), Inches(0.28))
    sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor(*color); sp.line.fill.background()
    T(s, pt, Inches(0.95), ty, Inches(6.5), Inches(0.38),
      size=13, color=WHITE)

T(s, "VivriMarket  —  2026", Inches(0.5), Inches(7.1), Inches(7), Inches(0.3),
  size=11, color=(60,120,70))


# ── Save ──────────────────────────────────────────────────────────────────────
out = "VivriMarket_Parcours_Client.pptx"
prs.save(out)
print(f"OK  Fichier genere : {out}")

# Cleanup temp images
import shutil
shutil.rmtree(TMP, ignore_errors=True)
print("   Temp images supprimees.")
