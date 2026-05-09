"""
VivriMarket_Nouvelles_Fonctionnalites.pptx
Présentation de mise en production — 4 nouvelles fonctionnalités
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x0A, 0x0A, 0x0A)

BG_NAVY = RGBColor(0x0D, 0x14, 0x1E)
BG_DARK = RGBColor(0x0D, 0x1F, 0x12)
BG_CHAT = RGBColor(0x07, 0x14, 0x28)
BG_STAT = RGBColor(0x0A, 0x18, 0x0A)
BG_SPON = RGBColor(0x1A, 0x10, 0x07)
BG_PAY  = RGBColor(0x12, 0x0A, 0x2A)

GREEN   = RGBColor(0x16, 0xA3, 0x4A)
LIME    = RGBColor(0x84, 0xCC, 0x16)
BLUE    = RGBColor(0x38, 0xBD, 0xF8)
VIOLET  = RGBColor(0xA7, 0x8B, 0xFA)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
YELLOW  = RGBColor(0xFD, 0xE6, 0x8A)
GOLD    = RGBColor(0xFB, 0xBF, 0x24)
ORANGE  = RGBColor(0xFB, 0x92, 0x3C)
RED     = RGBColor(0xF8, 0x71, 0x71)
TEAL    = RGBColor(0x2D, 0xD4, 0xBF)
CYAN    = RGBColor(0x22, 0xD3, 0xEE)
PINK    = RGBColor(0xF4, 0x72, 0xB6)

TXT_DIM = RGBColor(0x71, 0x77, 0x90)
TXT_MID = RGBColor(0xB0, 0xBA, 0xCC)
CODE_BG = RGBColor(0x0D, 0x13, 0x21)
CODE_FG = RGBColor(0x7D, 0xD3, 0xFC)

W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def oval(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(9, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def code_block(slide, lines, l, t, w, h, accent=BLUE):
    rect(slide, l, t, w, h, CODE_BG)
    rect(slide, l, t, Inches(0.07), h, accent)
    box = slide.shapes.add_textbox(l + Inches(0.2), t + Inches(0.15),
                                   w - Inches(0.3), h - Inches(0.28))
    tf = box.text_frame
    tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(12)
        r.font.name = "Consolas"
        if line.strip().startswith("#") or line.strip().startswith("//"):
            r.font.color.rgb = RGBColor(0x6A, 0x99, 0x6A)
        elif line.strip().startswith("⭐") or line.strip().startswith("✅"):
            r.font.color.rgb = GOLD
        else:
            r.font.color.rgb = CODE_FG


def top_bar(slide, color, h=Inches(0.1)):
    rect(slide, 0, 0, W, h, color)


def bottom_bar(slide, color, h=Inches(0.1)):
    rect(slide, 0, H - h, W, h, color)


def section_header(slide, num, title, accent, bg):
    rect(slide, 0, 0, W, Inches(0.1), accent)
    pill = oval(slide, Inches(0.7), Inches(0.28), Inches(0.68), Inches(0.68), accent)
    tf = pill.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = BLACK
    txt(slide, title, Inches(1.55), Inches(0.28), Inches(11.5), Inches(0.7),
        size=26, bold=True, color=accent)
    rect(slide, Inches(0.7), Inches(1.08), Inches(12.2), Inches(0.03), accent)


def kpi_card(slide, l, t, w, h, icon, value, label, accent, bg):
    rect(slide, l, t, w, h, bg)
    rect(slide, l, t, Inches(0.14), h, accent)
    txt(slide, icon,  l + Inches(0.28), t + Inches(0.18), Inches(0.8), Inches(0.7),
        size=26, align=PP_ALIGN.CENTER)
    txt(slide, value, l + Inches(1.1), t + Inches(0.12), w - Inches(1.25), Inches(0.55),
        size=22, bold=True, color=accent)
    txt(slide, label, l + Inches(1.1), t + Inches(0.65), w - Inches(1.25), Inches(0.4),
        size=11, color=TXT_MID)


def feature_tag(slide, label, l, t, color):
    w, h = Inches(2.5), Inches(0.42)
    rect(slide, l, t, w, h, color)
    txt(slide, label, l, t, w, h, size=12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITRE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_NAVY)

# Vertical bands right
for i, c in enumerate([
    RGBColor(0x16,0x42,0x2A), RGBColor(0x1E,0x2A,0x4A),
    RGBColor(0x2A,0x1A,0x4A), RGBColor(0x1A,0x2A,0x10),
]):
    rect(s, W - Inches(3.8) + i * Inches(1.0), 0, Inches(0.85), H, c)

# Left accent
rect(s, 0, 0, Inches(0.55), H, GREEN)
rect(s, Inches(0.55), 0, Inches(0.07), H, RGBColor(0x16,0x85,0x3A))

# Logo
oval(s, Inches(1.1), Inches(1.3), Inches(1.3), Inches(1.3), GREEN)
txt(s, "VM", Inches(1.1), Inches(1.3), Inches(1.3), Inches(1.3),
    size=34, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

txt(s, "VivriMarket", Inches(2.8), Inches(1.25), Inches(9.5), Inches(1.2),
    size=52, bold=True, color=WHITE)
txt(s, "AgriExchange Platform", Inches(2.85), Inches(2.5), Inches(8.5), Inches(0.65),
    size=22, color=GREEN)

# Badge titre
rect(s, Inches(2.8), Inches(3.4), Inches(9.2), Inches(0.78), GREEN)
txt(s, "  Présentation des nouvelles fonctionnalités — GO Production  ",
    Inches(2.8), Inches(3.4), Inches(9.2), Inches(0.78),
    size=18, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# 4 feature pills
features = [
    ("💬 Chat temps réel",        BLUE),
    ("📊 Dashboard stats",        TEAL),
    ("⭐ Sponsoring gratuit",      AMBER),
    ("💳 Sponsoring payant",      VIOLET),
]
for i, (label, color) in enumerate(features):
    lx = Inches(2.8) + i * Inches(2.35)
    rect(s, lx, Inches(4.55), Inches(2.2), Inches(0.52), color)
    txt(s, label, lx, Inches(4.55), Inches(2.2), Inches(0.52),
        size=12, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

txt(s, "Mai 2026  |  Mise en production VivriMarket",
    Inches(2.8), Inches(5.5), Inches(9), Inches(0.4),
    size=12, color=TXT_DIM)

bottom_bar(s, GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
top_bar(s, GREEN)
bottom_bar(s, GREEN)

txt(s, "4 nouvelles fonctionnalités", Inches(0.7), Inches(0.3), Inches(12), Inches(0.8),
    size=30, bold=True, color=GREEN)

cards = [
    ("01", "💬", "Chat temps réel",
     "Messagerie instantanée agriculteur ↔ acheteur via Socket.IO.\nBadge non-lus dans la barre de navigation.",
     BLUE, RGBColor(0x04,0x14,0x28)),
    ("02", "📊", "Dashboard stats agriculteur",
     "Tableau de bord enrichi : vues, contacts, note, chiffre d'affaires.\nGraphique de répartition par catégorie.",
     TEAL, RGBColor(0x04,0x1E,0x1A)),
    ("03", "⭐", "Produits sponsorisés (gratuit)",
     "Mise en avant des produits selon le forfait (BLEU=1, GOLD=3, PLATINUM=5).\nCarousel homepage + badge doré.",
     AMBER, RGBColor(0x28,0x14,0x02)),
    ("04", "💳", "Sponsoring payant — 5 000 FCFA",
     "Quand la limite gratuite est atteinte, payer 14 jours de visibilité.\nFlux CinetPay + expiration automatique.",
     VIOLET, RGBColor(0x18,0x0C,0x2A)),
]

for i, (num, icon, title, desc, accent, bg) in enumerate(cards):
    col = i % 2
    row = i // 2
    lx = Inches(0.7) + col * Inches(6.35)
    ty = Inches(1.25) + row * Inches(2.8)

    rect(s, lx, ty, Inches(6.05), Inches(2.55), bg)
    rect(s, lx, ty, Inches(0.14), Inches(2.55), accent)

    oval(s, lx + Inches(0.28), ty + Inches(0.28), Inches(0.6), Inches(0.6), accent)
    tf = s.shapes[-1].text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(13)
    r.font.bold = True
    r.font.color.rgb = BLACK

    txt(s, f"{icon}  {title}", lx + Inches(1.05), ty + Inches(0.25),
        Inches(4.8), Inches(0.55), size=17, bold=True, color=WHITE)
    txt(s, desc, lx + Inches(1.05), ty + Inches(0.88),
        Inches(4.8), Inches(1.4), size=12, color=TXT_MID)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CHAT TEMPS RÉEL : APERÇU
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_CHAT)
section_header(s, "01", "Chat temps réel — Fonctionnement", BLUE, BG_CHAT)

# Left column: what it does
txt(s, "Ce que l'utilisateur voit", Inches(0.7), Inches(1.2), Inches(5.8), Inches(0.5),
    size=14, bold=True, color=BLUE)

items_user = [
    ("💬", "Icône Messages dans la navbar avec badge rouge (non-lus)"),
    ("📋", "Liste des conversations avec nom, avatar et horodatage"),
    ("✉️", "Zone de chat : bulles vertes (moi) / blanches (autre)"),
    ("🔗", "Bouton « Contacter le vendeur » sur chaque fiche produit"),
    ("🔴", "Badge se met à zéro quand la conversation est ouverte"),
]
for i, (icon, desc) in enumerate(items_user):
    ty = Inches(1.82) + i * Inches(0.94)
    rect(s, Inches(0.7), ty, Inches(5.8), Inches(0.8), RGBColor(0x06,0x18,0x32))
    rect(s, Inches(0.7), ty, Inches(0.12), Inches(0.8), BLUE)
    txt(s, icon, Inches(0.95), ty + Inches(0.1), Inches(0.6), Inches(0.6), size=18)
    txt(s, desc, Inches(1.6), ty + Inches(0.15), Inches(4.75), Inches(0.52),
        size=12, color=TXT_MID)

# Right column: mock UI
rect(s, Inches(7.0), Inches(1.2), Inches(5.9), Inches(5.8), RGBColor(0x04,0x10,0x20))
rect(s, Inches(7.0), Inches(1.2), Inches(5.9), Inches(0.55), RGBColor(0x06,0x20,0x3A))
txt(s, "💬  Messages", Inches(7.1), Inches(1.22), Inches(3), Inches(0.5),
    size=13, bold=True, color=BLUE)

# Badge example in nav
rect(s, Inches(11.9), Inches(1.25), Inches(0.7), Inches(0.42), RGBColor(0xDC,0x26,0x26))
txt(s, "3", Inches(11.9), Inches(1.25), Inches(0.7), Inches(0.42),
    size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Sidebar
rect(s, Inches(7.0), Inches(1.75), Inches(2.2), Inches(5.25), RGBColor(0x06,0x14,0x28))
for i, (name, preview) in enumerate([
    ("Amadou D.", "Bonjour, disponible ?"),
    ("Marie K.",  "Le prix du kilo ?"),
    ("Jean-P.",   "Je prends 50 kg"),
]):
    ty = Inches(1.85) + i * Inches(1.05)
    rect(s, Inches(7.05), ty, Inches(2.1), Inches(0.88), RGBColor(0x08,0x1C,0x34))
    oval(s, Inches(7.1), ty + Inches(0.15), Inches(0.55), Inches(0.55), BLUE)
    txt(s, name[0], Inches(7.1), ty + Inches(0.15), Inches(0.55), Inches(0.55),
        size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txt(s, name, Inches(7.72), ty + Inches(0.1), Inches(1.35), Inches(0.35),
        size=11, bold=True, color=WHITE)
    txt(s, preview, Inches(7.72), ty + Inches(0.46), Inches(1.35), Inches(0.32),
        size=9, color=TXT_DIM)

# Chat bubbles
rect(s, Inches(9.25), Inches(1.75), Inches(3.6), Inches(5.25), RGBColor(0x04,0x10,0x1C))
bubbles = [
    (True,  "Bonjour ! Vous avez des tomates ?",  Inches(2.0), BLUE),
    (False, "Oui, 200 kg disponibles.",             Inches(1.7), GREEN),
    (True,  "Quel est le prix au kilo ?",          Inches(1.8), BLUE),
    (False, "350 FCFA/kg franco",                   Inches(1.5), GREEN),
]
for i, (mine, msg, bw, bc) in enumerate(bubbles):
    ty = Inches(2.0) + i * Inches(0.85)
    bx = Inches(12.75) - bw if mine else Inches(9.35)
    rect(s, bx, ty, bw, Inches(0.65), bc if mine else RGBColor(0x1A,0x2A,0x3A))
    txt(s, msg, bx + Inches(0.1), ty + Inches(0.1), bw - Inches(0.2), Inches(0.45),
        size=10, color=WHITE if mine else TXT_MID)

# Input bar
rect(s, Inches(9.25), Inches(6.55), Inches(3.6), Inches(0.45),
     RGBColor(0x06,0x18,0x30))
txt(s, "Écrire un message…", Inches(9.35), Inches(6.58), Inches(2.8), Inches(0.38),
    size=10, color=TXT_DIM, italic=True)
oval(s, Inches(12.45), Inches(6.55), Inches(0.38), Inches(0.38), GREEN)
txt(s, "→", Inches(12.45), Inches(6.55), Inches(0.38), Inches(0.38),
    size=11, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

bottom_bar(s, BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CHAT TEMPS RÉEL : ARCHITECTURE TECHNIQUE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_CHAT)
section_header(s, "01", "Chat temps réel — Architecture technique", BLUE, BG_CHAT)

# Stack
txt(s, "Stack technique", Inches(0.7), Inches(1.2), Inches(5.8), Inches(0.5),
    size=14, bold=True, color=BLUE)

stack_items = [
    ("Socket.IO",  "WebSocket bidirectionnel — connexion persistante",  CYAN),
    ("MySQL",      "Table messages : sender_id, receiver_id, texte, lu", BLUE),
    ("JWT Auth",   "Token dans socket.handshake.auth pour sécuriser",    VIOLET),
    ("Salle perso","Chaque user rejoint socket.join(userId)",            TEAL),
]
for i, (tech, desc, color) in enumerate(stack_items):
    ty = Inches(1.82) + i * Inches(1.3)
    rect(s, Inches(0.7), ty, Inches(5.8), Inches(1.1), RGBColor(0x06,0x14,0x28))
    rect(s, Inches(0.7), ty, Inches(0.12), Inches(1.1), color)
    oval(s, Inches(0.92), ty + Inches(0.2), Inches(0.68), Inches(0.68), color)
    txt(s, tech, Inches(0.92), ty + Inches(0.2), Inches(0.68), Inches(0.68),
        size=8, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txt(s, tech, Inches(1.72), ty + Inches(0.14), Inches(4.6), Inches(0.45),
        size=13, bold=True, color=WHITE)
    txt(s, desc, Inches(1.72), ty + Inches(0.6), Inches(4.6), Inches(0.38),
        size=11, color=TXT_MID)

# Code snippet right
txt(s, "Flux d'événements Socket.IO", Inches(7.0), Inches(1.2), Inches(6.1), Inches(0.5),
    size=14, bold=True, color=BLUE)
code_block(s, [
    "// Agriculteur envoie un message",
    "socket.emit('send_message', {",
    "  receiverId: acheteurId,",
    "  texte: 'Oui 200kg disponibles'",
    "});",
    "",
    "// Serveur relaie aux deux parties",
    "io.to(senderId).to(receiverId)",
    "  .emit('new_message', message);",
    "",
    "// NavBar : incrément badge",
    "socket.on('new_message', () => {",
    "  if (path !== '/messages')",
    "    setUnreadCount(n => n + 1);",
    "});",
], Inches(7.0), Inches(1.82), Inches(6.1), Inches(4.9), BLUE)

# Files modified
txt(s, "Fichiers modifiés / créés",
    Inches(0.7), Inches(6.85), Inches(12.5), Inches(0.35),
    size=11, bold=True, color=TXT_DIM)
for i, (fname, color) in enumerate([
    ("SocketContext.jsx", VIOLET), ("MessagesPage.jsx + .css", BLUE),
    ("NavBar.jsx + .css", CYAN), ("ProductDetail.jsx + .css", TEAL),
    ("server.js (Socket.IO)", GREEN),
]):
    lx = Inches(0.7) + i * Inches(2.52)
    rect(s, lx, Inches(7.15), Inches(2.38), Inches(0.28), color)
    txt(s, fname, lx, Inches(7.15), Inches(2.38), Inches(0.28),
        size=9, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

bottom_bar(s, BLUE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DASHBOARD STATS AGRICULTEUR : APERÇU
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_STAT)
section_header(s, "02", "Dashboard stats agriculteur — Vue d'ensemble", TEAL, BG_STAT)

txt(s, "Nouvelles métriques de performance ajoutées au dashboard agriculteur",
    Inches(0.7), Inches(1.18), Inches(12.5), Inches(0.42),
    size=13, color=TXT_MID)

# 4 KPI cards
kpis = [
    ("👁", "1 248", "Vues produits",       BLUE,   RGBColor(0x04,0x14,0x28)),
    ("💬", "67",    "Contacts reçus",       TEAL,   RGBColor(0x04,0x1A,0x18)),
    ("⭐", "4.7/5", "Note moyenne",         GOLD,   RGBColor(0x28,0x1C,0x04)),
    ("💰", "342K",  "CA estimé (FCFA)",     GREEN,  RGBColor(0x04,0x1A,0x08)),
]
for i, (icon, val, label, accent, bg) in enumerate(kpis):
    kpi_card(s, Inches(0.7) + i * Inches(3.15), Inches(1.75),
             Inches(2.95), Inches(1.35), icon, val, label, accent, bg)

# Stock chips
txt(s, "État du stock en temps réel", Inches(0.7), Inches(3.28), Inches(12), Inches(0.4),
    size=13, bold=True, color=WHITE)

chips = [
    ("Tomates — 120 kg", GREEN,  RGBColor(0x06,0x20,0x0C)),
    ("Oignons — 8 kg",   AMBER,  RGBColor(0x28,0x14,0x02)),
    ("Maïs — 0 kg",      RED,    RGBColor(0x28,0x04,0x04)),
    ("Piment — 45 kg",   GREEN,  RGBColor(0x06,0x20,0x0C)),
    ("Aubergine — 5 kg", AMBER,  RGBColor(0x28,0x14,0x02)),
]
for i, (label, color, bg) in enumerate(chips):
    lx = Inches(0.7) + i * Inches(2.52)
    rect(s, lx, Inches(3.78), Inches(2.38), Inches(0.52), bg)
    rect(s, lx, Inches(3.78), Inches(0.12), Inches(0.52), color)
    txt(s, label, lx + Inches(0.22), Inches(3.82), Inches(2.1), Inches(0.38),
        size=11, color=WHITE)

# Category bar chart
txt(s, "Répartition par catégorie", Inches(0.7), Inches(4.5), Inches(6), Inches(0.4),
    size=13, bold=True, color=WHITE)
bars = [
    ("Légumes",  78, TEAL),
    ("Tubercules", 55, GREEN),
    ("Fruits",   38, LIME),
    ("Céréales", 22, AMBER),
    ("Épices",   15, ORANGE),
]
for i, (cat, pct, color) in enumerate(bars):
    ty = Inches(5.0) + i * Inches(0.42)
    txt(s, cat, Inches(0.7), ty, Inches(1.6), Inches(0.38), size=10, color=TXT_MID)
    total_w = Inches(5.0)
    bar_w = total_w * pct / 100
    rect(s, Inches(2.4), ty + Inches(0.06), bar_w, Inches(0.28), color)
    txt(s, f"{pct}%", Inches(2.4) + bar_w + Inches(0.12), ty, Inches(0.7), Inches(0.38),
        size=10, color=color)

# Star rating
txt(s, "Note clients", Inches(7.5), Inches(4.5), Inches(5.5), Inches(0.4),
    size=13, bold=True, color=WHITE)
for i in range(5):
    star_color = GOLD if i < 4 else TXT_DIM
    oval(s, Inches(7.5) + i * Inches(0.72), Inches(5.0),
         Inches(0.6), Inches(0.6), star_color)
    txt(s, "★", Inches(7.5) + i * Inches(0.72), Inches(5.0), Inches(0.6), Inches(0.6),
        size=18, color=BLACK, align=PP_ALIGN.CENTER)
txt(s, "4.7 / 5  (12 avis)", Inches(11.2), Inches(5.05), Inches(2), Inches(0.42),
    size=12, color=GOLD)

bottom_bar(s, TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DASHBOARD STATS : CÔTÉ SERVEUR
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_STAT)
section_header(s, "02", "Dashboard stats agriculteur — Backend", TEAL, BG_STAT)

txt(s, "Endpoint  GET /api/v1/products/my-stats  (agriculteur authentifié)",
    Inches(0.7), Inches(1.18), Inches(12.5), Inches(0.42),
    size=13, color=TEAL, italic=True)

code_block(s, [
    "// mysqlProductRepository.js  —  getSellerStats(sellerId)",
    "// 6 requêtes agrégées en parallèle :",
    "",
    "  1. COUNT(*)  FROM products WHERE seller_id = ?  → nb_produits",
    "  2. SUM(views) FROM product_views WHERE product_id IN (mes IDs) → vues",
    "  3. COUNT(*)  FROM messages WHERE receiver_id = ?   → contacts",
    "  4. AVG(note), COUNT(*) FROM seller_reviews WHERE seller_id = ?",
    "  5. SUM(amount) FROM product_payments WHERE seller_id = 'active'",
    "  6. GROUP BY categorie → répartition catégories",
    "",
    "// Réponse JSON : { stats: { produits, vues, contacts,",
    "//   noteMoyenne, nbAvis, revenuEstime, categories } }",
], Inches(0.7), Inches(1.72), Inches(8.5), Inches(4.6), TEAL)

# Right: files changed
txt(s, "Fichiers impactés", Inches(9.5), Inches(1.72), Inches(3.6), Inches(0.45),
    size=13, bold=True, color=WHITE)
files = [
    ("mysqlProductRepository.js", "+ getSellerStats()", TEAL),
    ("routes/products.js",         "GET /my-stats",       GREEN),
    ("FarmerDashboard.jsx",        "+ section Perf.",     CYAN),
    ("FarmerDashboard.css",        "+ KPI grid + bars",   BLUE),
]
for i, (file, change, color) in enumerate(files):
    ty = Inches(2.3) + i * Inches(1.25)
    rect(s, Inches(9.5), ty, Inches(3.6), Inches(1.05), RGBColor(0x04,0x1A,0x0C))
    rect(s, Inches(9.5), ty, Inches(0.12), Inches(1.05), color)
    txt(s, file, Inches(9.72), ty + Inches(0.1), Inches(3.28), Inches(0.42),
        size=11, bold=True, color=WHITE)
    txt(s, change, Inches(9.72), ty + Inches(0.56), Inches(3.28), Inches(0.38),
        size=10, color=color)

# Impact
rect(s, Inches(0.7), Inches(6.55), Inches(12.5), Inches(0.55),
     RGBColor(0x04,0x20,0x16))
rect(s, Inches(0.7), Inches(6.55), Inches(0.14), Inches(0.55), TEAL)
txt(s, "  Impact : L'agriculteur voit ses performances en temps réel — il peut adapter ses prix, stocks "
    "et stratégie de sponsoring.",
    Inches(1.0), Inches(6.6), Inches(12.1), Inches(0.45),
    size=12, color=TEAL)

bottom_bar(s, TEAL)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PRODUITS SPONSORISÉS (GRATUIT) : APERÇU
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_SPON)
section_header(s, "03", "Produits sponsorisés — Système gratuit", AMBER, BG_SPON)

# Left: how it works
txt(s, "Comment ça fonctionne", Inches(0.7), Inches(1.2), Inches(6.0), Inches(0.5),
    size=14, bold=True, color=AMBER)

steps = [
    ("1", "L'agriculteur clique « ☆ Sponsoriser » sur un produit dans Mes Produits"),
    ("2", "Le système vérifie sa limite selon son forfait actif"),
    ("3", "Si quota disponible → is_featured = 1 dans la base MySQL"),
    ("4", "Le produit apparaît dans le carousel ⭐ de la page d'accueil"),
    ("5", "Badge doré « Sponsorisé » affiché sur la fiche et en catégorie"),
    ("6", "L'agriculteur peut désactiver le sponsoring à tout moment"),
]
for i, (num, desc) in enumerate(steps):
    ty = Inches(1.82) + i * Inches(0.87)
    rect(s, Inches(0.7), ty, Inches(6.0), Inches(0.75), RGBColor(0x28,0x10,0x02))
    oval(s, Inches(0.82), ty + Inches(0.1), Inches(0.52), Inches(0.52), AMBER)
    tf = s.shapes[-1].text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = BLACK
    txt(s, desc, Inches(1.45), ty + Inches(0.14), Inches(5.15), Inches(0.48),
        size=11, color=TXT_MID)

# Right: Limits table + carousel preview
txt(s, "Limites par forfait", Inches(7.1), Inches(1.2), Inches(6.0), Inches(0.5),
    size=14, bold=True, color=AMBER)

plans = [
    ("BLEU",     "1 produit",  RGBColor(0x38,0xBD,0xF8), RGBColor(0x04,0x14,0x28)),
    ("GOLD",     "3 produits", GOLD,                      RGBColor(0x28,0x1C,0x04)),
    ("PLATINUM", "5 produits", TEAL,                      RGBColor(0x04,0x1A,0x18)),
]
for i, (plan, limit, color, bg) in enumerate(plans):
    ty = Inches(1.82) + i * Inches(0.82)
    rect(s, Inches(7.1), ty, Inches(5.9), Inches(0.7), bg)
    rect(s, Inches(7.1), ty, Inches(0.12), Inches(0.7), color)
    rect(s, Inches(7.3), ty + Inches(0.12), Inches(1.5), Inches(0.46), color)
    txt(s, plan, Inches(7.3), ty + Inches(0.12), Inches(1.5), Inches(0.46),
        size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txt(s, limit, Inches(9.0), ty + Inches(0.18), Inches(3.8), Inches(0.38),
        size=13, color=WHITE)

# Carousel mockup
txt(s, "Carousel homepage « Produits à la une »", Inches(7.1), Inches(4.4), Inches(6.0), Inches(0.45),
    size=13, bold=True, color=WHITE)
rect(s, Inches(7.1), Inches(4.95), Inches(5.9), Inches(2.2),
     RGBColor(0x1A,0x0E,0x04))
txt(s, "‹", Inches(7.1), Inches(5.6), Inches(0.5), Inches(0.9),
    size=28, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
for i, (name, price) in enumerate([("Tomates", "350 FCFA"), ("Maïs", "280 FCFA"), ("Oignons", "420 FCFA")]):
    lx = Inches(7.7) + i * Inches(1.7)
    rect(s, lx, Inches(5.05), Inches(1.55), Inches(2.0), RGBColor(0x28,0x14,0x04))
    oval(s, lx + Inches(0.22), Inches(5.12), Inches(1.12), Inches(0.92),
         RGBColor(0x40,0x20,0x06))
    txt(s, "🍅🌽🧅"[i], lx + Inches(0.22), Inches(5.12), Inches(1.12), Inches(0.92),
        size=24, align=PP_ALIGN.CENTER)
    rect(s, lx, Inches(6.1), Inches(1.55), Inches(0.28), AMBER)
    txt(s, "⭐ Sponsorisé", lx, Inches(6.1), Inches(1.55), Inches(0.28),
        size=8, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txt(s, name, lx + Inches(0.05), Inches(6.45), Inches(1.45), Inches(0.3),
        size=10, bold=True, color=WHITE)
    txt(s, price, lx + Inches(0.05), Inches(6.78), Inches(1.45), Inches(0.28),
        size=9, color=AMBER)
txt(s, "›", Inches(12.6), Inches(5.6), Inches(0.5), Inches(0.9),
    size=28, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

bottom_bar(s, AMBER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SPONSORING PAYANT : OFFRE & FLUX
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_PAY)
section_header(s, "04", "Sponsoring payant — 5 000 FCFA / 14 jours", VIOLET, BG_PAY)

# Offer box
rect(s, Inches(0.7), Inches(1.25), Inches(5.8), Inches(3.5),
     RGBColor(0x1A,0x0E,0x2A))
rect(s, Inches(0.7), Inches(1.25), Inches(5.8), Inches(0.14), VIOLET)

txt(s, "L'offre payante", Inches(0.9), Inches(1.45), Inches(5.4), Inches(0.55),
    size=16, bold=True, color=VIOLET)

# Price highlight
rect(s, Inches(0.9), Inches(2.1), Inches(2.2), Inches(1.2),
     RGBColor(0x2A,0x18,0x40))
txt(s, "5 000", Inches(0.9), Inches(2.1), Inches(2.2), Inches(0.85),
    size=38, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "FCFA", Inches(0.9), Inches(2.82), Inches(2.2), Inches(0.38),
    size=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

txt(s, "⏱  14 jours de visibilité", Inches(3.3), Inches(2.18), Inches(3.0), Inches(0.5),
    size=14, bold=True, color=YELLOW)
txt(s, "Activation immédiate\naprès confirmation du paiement",
    Inches(3.3), Inches(2.75), Inches(3.0), Inches(0.8),
    size=11, color=TXT_MID)

benefits = [
    "✅ Affichage prioritaire en page d'accueil",
    "✅ Badge doré « Sponsorisé » sur la fiche",
    "✅ Tête de liste dans sa catégorie",
    "✅ Expiration automatique (cron horaire)",
]
for i, b in enumerate(benefits):
    txt(s, b, Inches(0.9), Inches(3.5) + i * Inches(0.38), Inches(5.4), Inches(0.35),
        size=11, color=TXT_MID)

# Payment flow
txt(s, "Flux de paiement CinetPay", Inches(7.0), Inches(1.25), Inches(6.1), Inches(0.5),
    size=14, bold=True, color=VIOLET)

flow_steps = [
    ("1", "Limite gratuite atteinte",          "PUT /products/:id/sponsor → 403 limitReached", AMBER),
    ("2", "Modal SponsorPayModal s'ouvre",     "Affiche l'offre 5 000 FCFA / 14 jours",        VIOLET),
    ("3", "Clic « Payer »",                    "POST /products/:id/sponsor/initiate → CinetPay",  BLUE),
    ("4", "Redirection CinetPay",              "Mobile Money / carte bancaire acceptés",          CYAN),
    ("5", "Retour return_url",                 "?sponsor_tx=TX&sponsor_pid=PID (ou webhook)",    GREEN),
    ("6", "Vérification & activation",         "GET /sponsor/check → activateSponsor() MySQL",   TEAL),
]
for i, (num, step, detail, color) in enumerate(flow_steps):
    ty = Inches(1.88) + i * Inches(0.9)
    rect(s, Inches(7.0), ty, Inches(6.1), Inches(0.78),
         RGBColor(0x14,0x08,0x28) if i % 2 == 0 else RGBColor(0x10,0x06,0x22))
    oval(s, Inches(7.1), ty + Inches(0.12), Inches(0.52), Inches(0.52), color)
    tf = s.shapes[-1].text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = BLACK
    txt(s, step, Inches(7.74), ty + Inches(0.06), Inches(5.25), Inches(0.38),
        size=12, bold=True, color=WHITE)
    txt(s, detail, Inches(7.74), ty + Inches(0.44), Inches(5.25), Inches(0.3),
        size=9, color=color, italic=True)

# Cron note
rect(s, Inches(0.7), Inches(6.55), Inches(12.5), Inches(0.55),
     RGBColor(0x10,0x06,0x22))
rect(s, Inches(0.7), Inches(6.55), Inches(0.14), Inches(0.55), VIOLET)
txt(s, "  Cron horaire (server.js) : expireOldSponsorships() — remet is_featured=0 "
    "quand paid_sponsor_until < NOW(). Run immédiat au démarrage + toutes les heures.",
    Inches(1.0), Inches(6.6), Inches(12.1), Inches(0.45),
    size=11, color=TXT_MID)

bottom_bar(s, VIOLET)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ARCHITECTURE DONNÉES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_NAVY)
top_bar(s, GREEN)

txt(s, "Architecture données — Tables MySQL impactées", Inches(0.7), Inches(0.3),
    Inches(12), Inches(0.75), size=26, bold=True, color=GREEN)
rect(s, Inches(0.7), Inches(1.08), Inches(12.2), Inches(0.03), GREEN)

tables_info = [
    ("messages", [
        "id INT PK AUTO_INCREMENT",
        "sender_id INT FK → users",
        "receiver_id INT FK → users",
        "texte TEXT NOT NULL",
        "lu TINYINT DEFAULT 0",
        "produit_id INT NULL",
        "created_at DATETIME",
    ], BLUE, RGBColor(0x04,0x14,0x28)),
    ("products (colonnes ajoutées)", [
        "is_featured TINYINT(1) DEFAULT 0",
        "paid_sponsor_until DATETIME NULL",
        "— (colonne ajoutée via dbMigrations.js)",
        "— is_featured existait déjà",
    ], AMBER, RGBColor(0x28,0x14,0x02)),
    ("product_sponsor_payments", [
        "id INT PK AUTO_INCREMENT",
        "product_id INT FK → products",
        "seller_id INT FK → users",
        "transaction_id VARCHAR(100) UNIQUE",
        "amount DECIMAL(10,2) DEFAULT 5000",
        "status ENUM(pending,active,expired)",
        "start_date / end_date DATETIME",
    ], VIOLET, RGBColor(0x14,0x08,0x28)),
    ("seller_reviews (lectures)", [
        "note TINYINT (1-5)",
        "seller_id INT FK → users",
        "— Agrégation AVG(note) pour",
        "  le dashboard agriculteur",
    ], TEAL, RGBColor(0x04,0x1A,0x18)),
]

for i, (tname, columns, accent, bg) in enumerate(tables_info):
    col = i % 2
    row = i // 2
    lx = Inches(0.7) + col * Inches(6.35)
    ty = Inches(1.25) + row * Inches(2.9)

    rect(s, lx, ty, Inches(6.05), Inches(2.7), bg)
    rect(s, lx, ty, Inches(0.14), Inches(2.7), accent)
    rect(s, lx, ty, Inches(6.05), Inches(0.48), accent)
    txt(s, tname, lx + Inches(0.2), ty + Inches(0.06), Inches(5.7), Inches(0.38),
        size=13, bold=True, color=BLACK)
    for j, col_def in enumerate(columns[:6]):
        txt(s, f"  {col_def}", lx + Inches(0.2), ty + Inches(0.55) + j * Inches(0.35),
            Inches(5.7), Inches(0.32), size=10, color=TXT_MID if "—" in col_def else CODE_FG)

bottom_bar(s, GREEN)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CHECKLIST GO PRODUCTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK)
top_bar(s, GREEN)
bottom_bar(s, GREEN)

txt(s, "Checklist GO Production", Inches(0.7), Inches(0.3), Inches(12), Inches(0.75),
    size=30, bold=True, color=GREEN)
rect(s, Inches(0.7), Inches(1.08), Inches(12.2), Inches(0.03), GREEN)

checks = [
    # (label, detail, status_color, status_text)
    ("💬 Chat Socket.IO",       "Serveur + client testés en local, messages reçus en temps réel", GREEN,  "✅ Prêt"),
    ("📊 Dashboard stats",      "Endpoint /my-stats retourne les 6 métriques, UI affichée",       GREEN,  "✅ Prêt"),
    ("⭐ Sponsoring gratuit",    "Toggle is_featured, limits BLEU/GOLD/PLATINUM, carousel affiché",GREEN,  "✅ Prêt"),
    ("💳 Sponsoring payant",    "Flow CinetPay complet : initiate → retour → check → activation", GREEN,  "✅ Prêt"),
    ("🔁 Cron expiration",      "setInterval 1h + run immédiat au démarrage serveur",              GREEN,  "✅ Prêt"),
    ("🗄️ Migration MySQL",      "paid_sponsor_until ajouté via dbMigrations.js (idempotent)",     GREEN,  "✅ Prêt"),
    ("🔒 Auth JWT Socket",      "socket.handshake.auth.token vérifié avant tout événement",        GREEN,  "✅ Prêt"),
    ("📱 Responsive UI",        "Chat, dashboard, modal CinetPay — testés mobile 480px",           AMBER,  "⚠️ A tester"),
    ("🔑 Vars env production",  "CINETPAY_APIKEY, CINETPAY_SITE_ID, JWT_SECRET en prod",          AMBER,  "⚠️ Configurer"),
]

for i, (label, detail, status_color, status_text) in enumerate(checks):
    col = i % 3
    row = i // 3
    lx = Inches(0.7) + col * Inches(4.15)
    ty = Inches(1.25) + row * Inches(1.88)

    rect(s, lx, ty, Inches(3.95), Inches(1.68), RGBColor(0x0A,0x1E,0x10))
    rect(s, lx, ty, Inches(0.12), Inches(1.68), status_color)

    txt(s, label, lx + Inches(0.22), ty + Inches(0.12), Inches(3.55), Inches(0.48),
        size=12, bold=True, color=WHITE)
    txt(s, detail, lx + Inches(0.22), ty + Inches(0.64), Inches(3.55), Inches(0.65),
        size=9, color=TXT_MID)

    rect(s, lx + Inches(0.22), ty + Inches(1.3), Inches(2.3), Inches(0.28), status_color)
    txt(s, status_text, lx + Inches(0.22), ty + Inches(1.3), Inches(2.3), Inches(0.28),
        size=10, bold=True, color=BLACK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — RÉCAPITULATIF & DEMANDE DE GO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_NAVY)

# Colored right bands
for i, c in enumerate([GREEN, TEAL, AMBER, VIOLET, BLUE, LIME]):
    rect(s, W - Inches(3.2) + i * Inches(0.55), 0, Inches(0.52), H, c)

# Big GO circle
oval(s, Inches(1.0), Inches(1.5), Inches(2.2), Inches(2.2), GREEN)
txt(s, "GO", Inches(1.0), Inches(1.5), Inches(2.2), Inches(2.2),
    size=52, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
txt(s, "?", Inches(2.9), Inches(1.52), Inches(0.9), Inches(2.2),
    size=80, bold=True, color=GREEN)

txt(s, "4 fonctionnalités prêtes", Inches(4.3), Inches(1.5), Inches(8), Inches(0.9),
    size=36, bold=True, color=WHITE)
txt(s, "pour la mise en production", Inches(4.3), Inches(2.42), Inches(8), Inches(0.7),
    size=22, color=GREEN)

# 4 feature summary pills
features_summary = [
    ("💬", "Chat temps réel",        "Socket.IO + badge non-lus",           BLUE),
    ("📊", "Dashboard stats",         "6 métriques, graphiques, note",        TEAL),
    ("⭐", "Sponsoring gratuit",      "Carousel homepage, limites plan",       AMBER),
    ("💳", "Sponsoring payant",       "5 000 FCFA, CinetPay, cron expiry",    VIOLET),
]
for i, (icon, title, sub, color) in enumerate(features_summary):
    ty = Inches(3.55) + i * Inches(0.88)
    rect(s, Inches(1.0), ty, Inches(11.0), Inches(0.75),
         RGBColor(0x0A,0x18,0x10) if i % 2 == 0 else RGBColor(0x0A,0x10,0x18))
    rect(s, Inches(1.0), ty, Inches(0.12), Inches(0.75), color)
    txt(s, f"{icon}  {title}", Inches(1.22), ty + Inches(0.1),
        Inches(4.5), Inches(0.45), size=13, bold=True, color=WHITE)
    txt(s, sub, Inches(6.0), ty + Inches(0.18), Inches(5.8), Inches(0.38),
        size=11, color=color, italic=True)
    rect(s, Inches(11.5), ty + Inches(0.16), Inches(0.45), Inches(0.42), GREEN)
    txt(s, "✓", Inches(11.5), ty + Inches(0.16), Inches(0.45), Inches(0.42),
        size=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

txt(s, "VivriMarket / AgriExchange  —  Mai 2026",
    Inches(1.0), Inches(7.1), Inches(10), Inches(0.4),
    size=11, color=TXT_DIM)

top_bar(s, GREEN)
bottom_bar(s, GREEN)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "VivriMarket_Nouvelles_Fonctionnalites.pptx"
prs.save(out)
print(f"OK  Fichier sauvegardé : {out}")
