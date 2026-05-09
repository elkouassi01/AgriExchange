"""
AgriExchange_Installation_Guide.pptx — version colorée
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x0A, 0x0A, 0x0A)

# Slide backgrounds
BG_TITLE   = RGBColor(0x0D, 0x14, 0x1E)   # deep navy
BG_DARK1   = RGBColor(0x0D, 0x1F, 0x12)   # forest dark
BG_DARK2   = RGBColor(0x12, 0x0A, 0x2A)   # purple dark
BG_DARK3   = RGBColor(0x1A, 0x10, 0x07)   # amber dark
BG_DARK4   = RGBColor(0x07, 0x1A, 0x24)   # teal dark
BG_DARK5   = RGBColor(0x1A, 0x07, 0x14)   # pink dark
BG_DARK6   = RGBColor(0x0A, 0x18, 0x0A)   # deep green

# Accents
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
LIME       = RGBColor(0x84, 0xCC, 0x16)
BLUE       = RGBColor(0x38, 0xBD, 0xF8)
VIOLET     = RGBColor(0xA7, 0x8B, 0xFA)
ORANGE     = RGBColor(0xFB, 0x92, 0x3C)
YELLOW     = RGBColor(0xFB, 0xD5, 0x24)
PINK       = RGBColor(0xF4, 0x72, 0xB6)
CYAN       = RGBColor(0x22, 0xD3, 0xEE)
RED        = RGBColor(0xF8, 0x71, 0x71)
TEAL       = RGBColor(0x2D, 0xD4, 0xBF)

# Text shades
TXT_DIM    = RGBColor(0x71, 0x77, 0x90)
TXT_MID    = RGBColor(0xB0, 0xBA, 0xCC)
CODE_BG    = RGBColor(0x0D, 0x13, 0x21)
CODE_FG    = RGBColor(0x7D, 0xD3, 0xFC)

W = Inches(13.33)
H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color, line=False):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if not line:
        s.line.fill.background()
    return s


def oval(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(9, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return box


def code(slide, lines, l, t, w, h):
    rect(slide, l, t, w, h, CODE_BG)
    # left accent bar
    rect(slide, l, t, Inches(0.07), h, BLUE)
    box = slide.shapes.add_textbox(l + Inches(0.2), t + Inches(0.18),
                                   w - Inches(0.35), h - Inches(0.32))
    tf = box.text_frame
    tf.word_wrap = False
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.size = Pt(13)
        r.font.name = "Consolas"
        # color keywords
        if line.startswith("#"):
            r.font.color.rgb = RGBColor(0x6A, 0x99, 0x6A)
        elif any(line.strip().startswith(k) for k in ("npm","cd","node","python","mysql","git","Copy","cp","INSERT","CREATE","EXIT")):
            r.font.color.rgb = YELLOW
        else:
            r.font.color.rgb = CODE_FG


def side_bar(slide, color, width=Inches(0.4)):
    rect(slide, 0, 0, width, H, color)


def top_bar(slide, color, height=Inches(0.12)):
    rect(slide, 0, 0, W, height, color)


def bottom_bar(slide, color, height=Inches(0.1)):
    rect(slide, 0, H - height, W, height, color)


def slide_title(slide, step_num, title_text, accent):
    """Numbered pill + title line."""
    # accent top stripe
    rect(slide, 0, 0, W, Inches(0.1), accent)
    # step pill
    pill = oval(slide, Inches(0.7), Inches(0.3), Inches(0.65), Inches(0.65), accent)
    tf = pill.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = step_num
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = BLACK
    # title text
    txt(slide, title_text, Inches(1.55), Inches(0.3), Inches(11), Inches(0.7),
        size=26, bold=True, color=accent)
    # thin separator
    rect(slide, Inches(0.7), Inches(1.1), Inches(12), Inches(0.03), accent)


def tag_box(slide, label, l, t, w, h, bg_color, text_color=BLACK):
    rect(slide, l, t, w, h, bg_color)
    txt(slide, label, l + Inches(0.1), t + Inches(0.05),
        w - Inches(0.2), h - Inches(0.1),
        size=13, bold=True, color=text_color, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1  TITRE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_TITLE)

# Diagonal color bands (simulated with narrow rectangles)
for i, (clr, alpha_approx) in enumerate([
    (RGBColor(0x16,0x42,0x2A), Inches(0.3)),
    (RGBColor(0x1E,0x2A,0x4A), Inches(0.3)),
    (RGBColor(0x2A,0x1A,0x4A), Inches(0.3)),
]):
    rect(s, W - Inches(3.5) + i * Inches(1.1), 0, Inches(0.8), H, clr)

# Big vertical accent
rect(s, 0, 0, Inches(0.55), H, GREEN)
rect(s, Inches(0.55), 0, Inches(0.08), H, RGBColor(0x16,0x85,0x3A))

# Logo area
oval(s, Inches(1.1), Inches(1.4), Inches(1.2), Inches(1.2), GREEN)
txt(s, "AE", Inches(1.1), Inches(1.4), Inches(1.2), Inches(1.2),
    size=32, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# Title
txt(s, "AgriExchange", Inches(2.7), Inches(1.3), Inches(10), Inches(1.4),
    size=54, bold=True, color=WHITE)
txt(s, "VivriMarket", Inches(2.75), Inches(2.65), Inches(9), Inches(0.75),
    size=26, bold=False, color=GREEN)

# Subtitle badge
rect(s, Inches(2.7), Inches(3.55), Inches(8.5), Inches(0.75), GREEN)
txt(s, "  Guide d'installation sur une nouvelle machine  ",
    Inches(2.7), Inches(3.55), Inches(8.5), Inches(0.75),
    size=17, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

# Tech badges
for i, (label, color) in enumerate([
    ("Node.js", RGBColor(0x41,0x7E,0x38)),
    ("React",   RGBColor(0x09,0x7F,0xAA)),
    ("MySQL",   RGBColor(0x00,0x5C,0x84)),
    ("Vite",    RGBColor(0x74,0x6F,0xE5)),
]):
    bx = Inches(2.7) + i * Inches(2.15)
    rect(s, bx, Inches(4.6), Inches(1.95), Inches(0.52), color)
    txt(s, label, bx, Inches(4.6), Inches(1.95), Inches(0.52),
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "2026  |  Plateforme agricole full-stack",
    Inches(2.7), Inches(5.5), Inches(10), Inches(0.4),
    size=12, color=TXT_DIM)

bottom_bar(s, GREEN)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2  SOMMAIRE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK6)
top_bar(s, GREEN)
bottom_bar(s, GREEN)

txt(s, "Sommaire", Inches(0.7), Inches(0.3), Inches(12), Inches(0.75),
    size=28, bold=True, color=GREEN)

steps = [
    ("01", "Prérequis",                          BLUE),
    ("02", "Cloner le dépôt",                    CYAN),
    ("03", "Installer les dépendances",           LIME),
    ("04", "Variables d'environnement (.env)",    YELLOW),
    ("05", "Créer la base de données MySQL",      ORANGE),
    ("06", "Démarrer le projet",                  GREEN),
    ("07", "Créer un compte administrateur",      VIOLET),
    ("08", "Options avancées",                    PINK),
]

for i, (num, label, color) in enumerate(steps):
    col = i % 2
    row = i // 2
    lx = Inches(0.7) + col * Inches(6.35)
    ty = Inches(1.35) + row * Inches(1.35)

    rect(s, lx, ty, Inches(6.0), Inches(1.1), RGBColor(0x0F,0x22,0x14))
    rect(s, lx, ty, Inches(0.12), Inches(1.1), color)

    oval(s, lx + Inches(0.25), ty + Inches(0.28), Inches(0.55), Inches(0.55), color)
    tf = s.shapes[-1].text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = BLACK

    txt(s, label, lx + Inches(0.95), ty + Inches(0.3), Inches(5), Inches(0.5),
        size=15, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3  PRÉREQUIS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK4)

slide_title(s, "01", "Prérequis", BLUE)

prereqs = [
    ("Node.js  >= 18", "LTS recommandé  —  nodejs.org",
     RGBColor(0x41,0x7E,0x38), RGBColor(0x06,0x22,0x10)),
    ("MySQL  8.x", "dev.mysql.com/downloads/mysql/",
     RGBColor(0x00,0x5C,0x84), RGBColor(0x04,0x18,0x28)),
    ("Git", "git-scm.com  (version ≥ 2.40)",
     RGBColor(0xF0,0x52,0x33), RGBColor(0x28,0x0A,0x04)),
    ("npm  >= 9", "Inclus automatiquement avec Node.js",
     RGBColor(0x9B,0x59,0xB6), RGBColor(0x20,0x0A,0x2A)),
]

for i, (title, desc, accent, bg) in enumerate(prereqs):
    col = i % 2
    row = i // 2
    lx = Inches(0.7) + col * Inches(6.3)
    ty = Inches(1.35) + row * Inches(2.55)

    rect(s, lx, ty, Inches(6.0), Inches(2.25), bg)
    rect(s, lx, ty, Inches(0.16), Inches(2.25), accent)
    txt(s, title, lx + Inches(0.35), ty + Inches(0.35),
        Inches(5.5), Inches(0.6), size=19, bold=True, color=WHITE)
    txt(s, desc, lx + Inches(0.35), ty + Inches(1.0),
        Inches(5.5), Inches(0.5), size=13, color=TXT_MID)

    # Verification check badge
    oval(s, lx + Inches(4.85), ty + Inches(0.3), Inches(0.65), Inches(0.45), accent)
    txt(s, "REQ", lx + Inches(4.85), ty + Inches(0.3), Inches(0.65), Inches(0.45),
        size=9, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

bottom_bar(s, BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4  CLONER LE DÉPÔT
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK1)

slide_title(s, "02", "Cloner le dépôt", CYAN)

rect(s, Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.5), RGBColor(0x07,0x2A,0x2A))
txt(s, "  Ouvrez un terminal (PowerShell, Git Bash ou cmd) puis copiez ces commandes :",
    Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.5),
    size=14, color=CYAN)

code(s, [
    "# 1. Se placer dans le dossier de votre choix",
    "cd C:\\Projets",
    "",
    "# 2. Cloner le depot",
    "git clone https://github.com/<votre-org>/AgriExchange.git",
    "",
    "# 3. Entrer dans le dossier",
    "cd AgriExchange",
], Inches(0.7), Inches(1.95), Inches(11.9), Inches(3.6))

# Info box
rect(s, Inches(0.7), Inches(5.75), Inches(11.9), Inches(0.9),
     RGBColor(0x07,0x24,0x28))
rect(s, Inches(0.7), Inches(5.75), Inches(0.14), Inches(0.9), YELLOW)
txt(s, "  Remplacez <votre-org> par le compte GitHub réel du projet.\n"
    "  Si le dépôt est privé, assurez-vous d'avoir accès SSH ou un Personal Access Token.",
    Inches(1.0), Inches(5.78), Inches(11.4), Inches(0.85),
    size=12, italic=True, color=YELLOW)

bottom_bar(s, CYAN)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5  DÉPENDANCES
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK2)

slide_title(s, "03", "Installer les dépendances", LIME)

# Server block
rect(s, Inches(0.7), Inches(1.25), Inches(5.85), Inches(0.55),
     RGBColor(0x1E,0x38,0x09))
txt(s, "  Serveur  (Node.js / Express)",
    Inches(0.7), Inches(1.25), Inches(5.85), Inches(0.55),
    size=14, bold=True, color=LIME)
code(s, ["cd server", "npm install", "cd .."],
     Inches(0.7), Inches(1.85), Inches(5.85), Inches(1.55))

# Client block
rect(s, Inches(7.1), Inches(1.25), Inches(5.5), Inches(0.55),
     RGBColor(0x09,0x1E,0x38))
txt(s, "  Client  (React / Vite)",
    Inches(7.1), Inches(1.25), Inches(5.5), Inches(0.55),
    size=14, bold=True, color=BLUE)
code(s, ["cd client", "npm install", "cd .."],
     Inches(7.1), Inches(1.85), Inches(5.5), Inches(1.55))

# Arrow
txt(s, "→", Inches(6.25), Inches(2.4), Inches(0.6), Inches(0.6),
    size=30, bold=True, color=LIME, align=PP_ALIGN.CENTER)

# Result
rect(s, Inches(0.7), Inches(3.65), Inches(11.9), Inches(0.5),
     RGBColor(0x14,0x20,0x14))
txt(s, "  Résultat attendu :", Inches(0.7), Inches(3.65), Inches(11.9), Inches(0.5),
    size=14, bold=True, color=WHITE)
code(s, [
    "server/node_modules/   OK",
    "client/node_modules/   OK",
], Inches(0.7), Inches(4.2), Inches(11.9), Inches(1.15))

# Warning
rect(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(0.85),
     RGBColor(0x28,0x1C,0x04))
rect(s, Inches(0.7), Inches(5.6), Inches(0.14), Inches(0.85), ORANGE)
txt(s, "  Si vous voyez des npm WARN peer : c'est généralement ignorable.\n"
    "  En cas d'erreur EACCES (Linux/macOS) : sudo chown -R $(whoami) ~/.npm",
    Inches(1.0), Inches(5.63), Inches(11.3), Inches(0.8),
    size=12, color=ORANGE)

bottom_bar(s, LIME)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6  VARIABLES D'ENVIRONNEMENT
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK3)

slide_title(s, "04", "Variables d'environnement  (.env)", YELLOW)

rect(s, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.5),
     RGBColor(0x28,0x20,0x04))
txt(s, "  1)  Copier le fichier exemple :",
    Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.5),
    size=14, bold=True, color=YELLOW)
code(s, [
    "# Windows (PowerShell)",
    "Copy-Item server\\.env.example  server\\.env",
    "",
    "# Linux / macOS",
    "cp server/.env.example  server/.env",
], Inches(0.7), Inches(1.82), Inches(11.9), Inches(2.05))

txt(s, "  2)  Valeurs obligatoires à renseigner :",
    Inches(0.7), Inches(4.05), Inches(11.9), Inches(0.45),
    size=14, bold=True, color=WHITE)

vars_info = [
    ("MYSQL_PASSWORD",  "Mot de passe de votre compte MySQL root",  RED),
    ("JWT_SECRET",      "Chaîne aléatoire >= 32 caractères",         VIOLET),
    ("REFRESH_SECRET",  "Autre chaîne aléatoire >= 32 caractères",   VIOLET),
    ("MYSQL_DATABASE",  "Nom de la base  (défaut : agriexchange)",   CYAN),
]
for i, (var, desc, color) in enumerate(vars_info):
    lx = Inches(0.7) + (i % 2) * Inches(5.95)
    ty = Inches(4.6) + (i // 2) * Inches(0.95)
    rect(s, lx, ty, Inches(5.7), Inches(0.82), RGBColor(0x18,0x0E,0x02))
    rect(s, lx, ty, Inches(0.14), Inches(0.82), color)
    txt(s, var, lx + Inches(0.25), ty + Inches(0.07), Inches(5.3), Inches(0.4),
        size=13, bold=True, color=color)
    txt(s, desc, lx + Inches(0.25), ty + Inches(0.46), Inches(5.3), Inches(0.3),
        size=11, color=TXT_MID)

bottom_bar(s, YELLOW)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7  BASE DE DONNÉES MYSQL
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK5)

slide_title(s, "05", "Créer la base de données MySQL", ORANGE)

txt(s, "1)  Créer la base (une seule fois) :",
    Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.45),
    size=14, bold=True, color=WHITE)
code(s, [
    "mysql -u root -p",
    "CREATE DATABASE agriexchange CHARACTER SET utf8mb4 COLLATE utf8mb4_unicode_ci;",
    "EXIT;",
], Inches(0.7), Inches(1.77), Inches(11.9), Inches(1.55))

txt(s, "2)  Lancer la migration automatique :",
    Inches(0.7), Inches(3.5), Inches(11.9), Inches(0.45),
    size=14, bold=True, color=WHITE)
code(s, [
    "cd server",
    "node scripts/migrate.js",
    "",
    "# Toutes les tables sont creees automatiquement",
], Inches(0.7), Inches(4.02), Inches(11.9), Inches(1.65))

# Tables badge area
rect(s, Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.7),
     RGBColor(0x28,0x0E,0x1A))
rect(s, Inches(0.7), Inches(5.9), Inches(0.14), Inches(0.7), ORANGE)
tables = ["users", "produits", "abonnements", "contact_requests", "paiements"]
for i, tbl in enumerate(tables):
    lx = Inches(1.1) + i * Inches(2.25)
    oval(s, lx, Inches(6.0), Inches(1.9), Inches(0.42), RGBColor(0x40,0x1A,0x06))
    txt(s, tbl, lx, Inches(6.0), Inches(1.9), Inches(0.42),
        size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

bottom_bar(s, ORANGE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8  DÉMARRER LE PROJET
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK6)

slide_title(s, "06", "Démarrer le projet", GREEN)

# Terminal 1
rect(s, Inches(0.7), Inches(1.25), Inches(5.85), Inches(0.55),
     RGBColor(0x09,0x23,0x18))
txt(s, "  Terminal 1  —  Serveur API  (port 5000)",
    Inches(0.7), Inches(1.25), Inches(5.85), Inches(0.55),
    size=13, bold=True, color=GREEN)
code(s, [
    "cd server",
    "npm run dev",
    "",
    "# Sortie attendue :",
    "# MySQL pool ready",
    "# Serveur demarre sur :5000",
], Inches(0.7), Inches(1.85), Inches(5.85), Inches(2.75))

# Terminal 2
rect(s, Inches(7.15), Inches(1.25), Inches(5.5), Inches(0.55),
     RGBColor(0x04,0x14,0x26))
txt(s, "  Terminal 2  —  Client React  (port 5173)",
    Inches(7.15), Inches(1.25), Inches(5.5), Inches(0.55),
    size=13, bold=True, color=BLUE)
code(s, [
    "cd client",
    "npm run dev",
    "",
    "# Sortie attendue :",
    "# VITE v5.x  ready",
    "#   Local: http://localhost:5173/",
], Inches(7.15), Inches(1.85), Inches(5.5), Inches(2.75))

# CTA
rect(s, Inches(0.7), Inches(4.85), Inches(11.9), Inches(1.05), GREEN)
txt(s, "Ouvrez  http://localhost:5173  dans votre navigateur",
    Inches(0.7), Inches(4.85), Inches(11.9), Inches(1.05),
    size=20, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

txt(s, "Les deux serveurs doivent être démarrés en même temps.",
    Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.4),
    size=12, italic=True, color=TXT_DIM, align=PP_ALIGN.CENTER)

bottom_bar(s, GREEN)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9  COMPTE ADMIN
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK2)

slide_title(s, "07", "Créer un compte administrateur", VIOLET)

# Option A
rect(s, Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.5),
     RGBColor(0x18,0x0C,0x2A))
txt(s, "  Option A  —  Script Node.js  (recommandé)",
    Inches(0.7), Inches(1.25), Inches(11.9), Inches(0.5),
    size=14, bold=True, color=VIOLET)
code(s, [
    'cd server',
    'node -e "',
    "  const bcrypt = require('bcryptjs');",
    "  const { getMysqlPool } = require('./config/mysql');",
    "  const pool = getMysqlPool();",
    "  const hash = bcrypt.hashSync('Admin1234!', 12);",
    "  pool.query('INSERT INTO users (nom,email,motDePasse,role,estActif,isVerified) VALUES (?,?,?,?,?,?)',",
    "    ['Admin','admin@agri.ci',hash,'admin',1,1]);",
    '"',
], Inches(0.7), Inches(1.82), Inches(11.9), Inches(3.35))

# Option B
rect(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(0.5),
     RGBColor(0x18,0x0C,0x2A))
txt(s, "  Option B  —  Directement en SQL",
    Inches(0.7), Inches(5.4), Inches(11.9), Inches(0.5),
    size=14, bold=True, color=PINK)
code(s, [
    "INSERT INTO users (nom, email, motDePasse, role, estActif, isVerified)",
    "VALUES ('Admin', 'admin@agri.ci', '<hash_bcrypt>', 'admin', 1, 1);",
], Inches(0.7), Inches(5.97), Inches(11.9), Inches(1.1))

bottom_bar(s, VIOLET)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10  OPTIONS AVANCÉES
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK1)

slide_title(s, "08", "Options avancées", PINK)

options = [
    ("Cloudinary",  "Stockage images produits",
     ["CLD_CLOUD=your_cloud_name", "CLD_KEY=your_api_key", "CLD_SECRET=your_api_secret"],
     BLUE,   RGBColor(0x04,0x14,0x26)),
    ("CinetPay",   "Paiements mobiles (CI, SN...)",
     ["CINETPAY_APIKEY=your_api_key", "CINETPAY_SITE_ID=your_site_id"],
     YELLOW, RGBColor(0x26,0x1E,0x04)),
    ("WhatsApp",   "Notifications OTP via WhatsApp",
     ["WHATSAPP_ENABLED=true", "# + scanner le QR code au 1er demarrage"],
     GREEN,  RGBColor(0x04,0x1A,0x0A)),
]

for i, (name, desc, lines, accent, bg) in enumerate(options):
    ty = Inches(1.3) + i * Inches(1.95)
    rect(s, Inches(0.7), ty, Inches(11.9), Inches(1.75), bg)
    rect(s, Inches(0.7), ty, Inches(0.18), Inches(1.75), accent)

    # Left: name + desc
    txt(s, name, Inches(1.1), ty + Inches(0.2), Inches(3.5), Inches(0.6),
        size=18, bold=True, color=accent)
    txt(s, desc, Inches(1.1), ty + Inches(0.82), Inches(3.5), Inches(0.5),
        size=12, color=TXT_MID)

    # Right: code
    code(s, lines, Inches(4.8), ty + Inches(0.18),
         Inches(7.5), Inches(1.4))

bottom_bar(s, PINK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11  CHECKLIST
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_DARK4)
top_bar(s, TEAL)

txt(s, "Checklist — Tout est prêt ?", Inches(0.7), Inches(0.3), Inches(12), Inches(0.75),
    size=28, bold=True, color=TEAL)
rect(s, Inches(0.7), Inches(1.1), Inches(12), Inches(0.03), TEAL)

items = [
    ("Node.js + MySQL + Git installés",          BLUE),
    ("Dépôt cloné  →  cd AgriExchange",          CYAN),
    ("npm install  dans server/ et client/",      LIME),
    ("server/.env  créé et configuré",            YELLOW),
    ("Base agriexchange créée (MySQL)",           ORANGE),
    ("Migration lancée  →  migrate.js",           ORANGE),
    ("Serveur démarré  (port 5000)",              GREEN),
    ("Client démarré   (port 5173)",              GREEN),
    ("Compte admin créé",                         VIOLET),
]

colors_cycle = [BLUE, CYAN, LIME, YELLOW, ORANGE, ORANGE, GREEN, GREEN, VIOLET]

for i, (item, color) in enumerate(zip([x for x,_ in items], colors_cycle)):
    col = i % 3
    row = i // 3
    lx = Inches(0.7) + col * Inches(4.15)
    ty = Inches(1.35) + row * Inches(1.75)

    rect(s, lx, ty, Inches(3.9), Inches(1.55), RGBColor(0x06,0x1C,0x22))
    rect(s, lx, ty, Inches(0.14), Inches(1.55), color)

    # Tick circle
    oval(s, lx + Inches(0.25), ty + Inches(0.2), Inches(0.55), Inches(0.55), color)
    tf_shape = s.shapes[-1]
    tf = tf_shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(i + 1)
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = BLACK

    txt(s, item, lx + Inches(0.92), ty + Inches(0.2), Inches(2.8), Inches(0.75),
        size=12, color=WHITE)

    # "A faire" chip
    rect(s, lx + Inches(0.25), ty + Inches(1.05), Inches(1.3), Inches(0.35), color)
    txt(s, "A faire", lx + Inches(0.25), ty + Inches(1.05), Inches(1.3), Inches(0.35),
        size=10, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

bottom_bar(s, TEAL)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12  FIN
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
set_bg(s, BG_TITLE)

# Colorful vertical bands on right
band_colors = [GREEN, CYAN, LIME, YELLOW, ORANGE, VIOLET, PINK, BLUE, TEAL]
for i, c in enumerate(band_colors):
    rect(s, W - Inches(3.5) + i * Inches(0.4), 0, Inches(0.38), H, c)

# Big tick
oval(s, Inches(1.5), Inches(1.4), Inches(1.6), Inches(1.6), GREEN)
txt(s, "OK", Inches(1.5), Inches(1.4), Inches(1.6), Inches(1.6),
    size=38, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

txt(s, "Installation terminée !",
    Inches(3.7), Inches(1.5), Inches(8), Inches(1.0),
    size=40, bold=True, color=WHITE)

txt(s, "Votre plateforme AgriExchange est opérationnelle.",
    Inches(3.7), Inches(2.55), Inches(8), Inches(0.6),
    size=18, color=TXT_MID)

# Color strip
for i, (label, color) in enumerate([
    ("Produits", GREEN), ("Abonnements", YELLOW), ("Contacts", BLUE),
    ("Paiements", ORANGE), ("Admin", VIOLET),
]):
    lx = Inches(1.5) + i * Inches(2.1)
    rect(s, lx, Inches(3.5), Inches(1.95), Inches(0.6), color)
    txt(s, label, lx, Inches(3.5), Inches(1.95), Inches(0.6),
        size=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

txt(s, "Pour toute question, ouvrez une issue sur GitHub",
    Inches(1.5), Inches(4.5), Inches(9), Inches(0.5),
    size=14, color=TXT_MID)

txt(s, "AgriExchange / VivriMarket  —  2026",
    Inches(1.5), Inches(6.8), Inches(10), Inches(0.4),
    size=11, color=TXT_DIM)

top_bar(s, GREEN)
bottom_bar(s, GREEN)


# ── Save ─────────────────────────────────────────────────────────────────────
out = "AgriExchange_Installation_Guide.pptx"
prs.save(out)
print(f"OK  Fichier sauvegarde : {out}")
